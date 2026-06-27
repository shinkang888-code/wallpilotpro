#!/usr/bin/env python3
"""Generate WallPilot Pro IR / Business Plan PowerPoint (25 slides)."""
from __future__ import annotations

import os
from pathlib import Path

import matplotlib.pyplot as plt
import matplotlib.font_manager as fm
from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
from pptx.util import Inches, Pt

ROOT = Path(__file__).resolve().parents[1]
ASSETS = ROOT / "assets"
DOWNLOADS = Path.home() / "Downloads"
OUT_DIR = DOWNLOADS / "WallPilot_Pro_IR_2026-06-27_v2"
OUT_PPTX = OUT_DIR / "WallPilot_Pro_IR_Business_Plan_25slides_v2.pptx"
SCREENS = ASSETS / "ir-screens"

# Brand colors
NAVY = RGBColor(0x0A, 0x0A, 0x0A)
GOLD = RGBColor(0xD4, 0xAF, 0x37)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
GRAY = RGBColor(0x9C, 0xA3, 0xAF)
CYAN = RGBColor(0x22, 0xD3, 0xEE)
GREEN = RGBColor(0x34, 0xD3, 0x99)
RED = RGBColor(0xF8, 0x71, 0x71)

# Korean font fallback
for f in ["Malgun Gothic", "AppleGothic", "NanumGothic", "DejaVu Sans"]:
    if any(f.lower() in x.name.lower() for x in fm.fontManager.ttflist):
        plt.rcParams["font.family"] = f
        break
plt.rcParams["axes.unicode_minus"] = False

CHART_STYLE = {"figure.facecolor": "#0a0a0a", "axes.facecolor": "#111827", "axes.edgecolor": "#374151", "axes.labelcolor": "#e5e7eb", "xtick.color": "#9ca3af", "ytick.color": "#9ca3af", "text.color": "#f3f4f6", "grid.color": "#1f2937"}


def ensure_assets() -> None:
    ASSETS.mkdir(parents=True, exist_ok=True)
    OUT_DIR.mkdir(parents=True, exist_ok=True)


def chart_mrr_scenarios() -> Path:
    months = ["M1", "M2", "M3", "M6", "M12"]
    conservative = [0.3, 0.8, 1.7, 4.0, 8.0]
    base = [0.8, 2.5, 4.6, 12.0, 28.0]
    optimistic = [1.5, 5.0, 11.4, 25.0, 55.0]

    fig, ax = plt.subplots(figsize=(10, 5.5))
    for k, v in CHART_STYLE.items():
        setattr(ax, k.split(".")[-1] if "." not in k else k, v) if False else None
    fig.patch.set_facecolor("#0a0a0a")
    ax.set_facecolor("#111827")
    ax.plot(months, conservative, "o-", color="#94a3b8", linewidth=2.5, label="보수 (₩M)")
    ax.plot(months, base, "o-", color="#d4af37", linewidth=3, label="기준 (₩M)")
    ax.plot(months, optimistic, "o-", color="#22d3ee", linewidth=2.5, label="낙관 (₩M)")
    ax.set_title("MRR 성장 시나리오 (KR+US)", fontsize=14, color="#f3f4f6", pad=12)
    ax.set_ylabel("MRR (백만 원)", color="#e5e7eb")
    ax.grid(True, alpha=0.3)
    ax.legend(facecolor="#1f2937", edgecolor="#374151", labelcolor="#f3f4f6")
    p = ASSETS / "chart_mrr_scenarios.png"
    fig.tight_layout()
    fig.savefig(p, dpi=180, facecolor="#0a0a0a")
    plt.close(fig)
    return p


def chart_tier_mix() -> Path:
    labels = ["Day\n60%", "Premium\n30%", "Elite\n10%"]
    sizes = [60, 30, 10]
    colors = ["#d4af37", "#22d3ee", "#34d399"]
    fig, ax = plt.subplots(figsize=(8, 5.5))
    fig.patch.set_facecolor("#0a0a0a")
    wedges, texts, autotexts = ax.pie(
        sizes, labels=labels, colors=colors, autopct="%1.0f%%", startangle=90,
        textprops={"color": "#f3f4f6", "fontsize": 11},
    )
    for t in autotexts:
        t.set_color("#0a0a0a")
        t.set_fontweight("bold")
    ax.set_title("유료 구독 Tier Mix (Blended ARPU ₩56.2K)", color="#f3f4f6", fontsize=13)
    p = ASSETS / "chart_tier_mix.png"
    fig.savefig(p, dpi=180, facecolor="#0a0a0a")
    plt.close(fig)
    return p


def chart_revenue_streams() -> Path:
    streams = ["B2C\n구독", "B2B\nAPI", "PoC/\nSetup", "US\nStripe"]
    m3_base = [4.0, 0, 0.675, 0.97]
    m12_base = [22, 8, 2, 6]
    x = range(len(streams))
    w = 0.35
    fig, ax = plt.subplots(figsize=(10, 5.5))
    fig.patch.set_facecolor("#0a0a0a")
    ax.set_facecolor("#111827")
    ax.bar([i - w / 2 for i in x], m3_base, w, label="Month 3 (₩M)", color="#d4af37")
    ax.bar([i + w / 2 for i in x], m12_base, w, label="Month 12 (₩M)", color="#22d3ee")
    ax.set_xticks(list(x))
    ax.set_xticklabels(streams, color="#e5e7eb")
    ax.set_title("수익원 구성 (기준 시나리오)", color="#f3f4f6", fontsize=13)
    ax.set_ylabel("₩M", color="#e5e7eb")
    ax.legend(facecolor="#1f2937", labelcolor="#f3f4f6")
    ax.grid(axis="y", alpha=0.25)
    p = ASSETS / "chart_revenue_streams.png"
    fig.tight_layout()
    fig.savefig(p, dpi=180, facecolor="#0a0a0a")
    plt.close(fig)
    return p


def chart_pnl() -> Path:
    months = ["Month 1", "Month 2", "Month 3"]
    mrr = [0.8, 2.5, 4.6]
    marketing = [-1.0, -1.5, -2.0]
    infra = [-0.3, -0.4, -0.5]
    net = [m + k + i for m, k, i in zip(mrr, marketing, infra)]
    fig, ax = plt.subplots(figsize=(10, 5.5))
    fig.patch.set_facecolor("#0a0a0a")
    ax.set_facecolor("#111827")
    ax.bar(months, mrr, label="MRR", color="#34d399")
    ax.bar(months, marketing, label="Marketing", color="#f87171")
    ax.bar(months, infra, bottom=marketing, label="Infra", color="#94a3b8")
    ax.plot(months, net, "o--", color="#d4af37", linewidth=2.5, label="Net (rough)")
    ax.axhline(0, color="#374151", linewidth=0.8)
    ax.set_title("90일 P&L 스케치 (기준, ₩M)", color="#f3f4f6", fontsize=13)
    ax.legend(facecolor="#1f2937", labelcolor="#f3f4f6")
    ax.grid(axis="y", alpha=0.25)
    p = ASSETS / "chart_pnl.png"
    fig.tight_layout()
    fig.savefig(p, dpi=180, facecolor="#0a0a0a")
    plt.close(fig)
    return p


def chart_market_tam() -> Path:
    labels = ["TAM\n글로번\n리테일\n퀀트·AI", "SAM\nKR+US\n액티브\n트레이더", "SOM\nY3\nWallPilot"]
    values = [120, 18, 2.4]
    fig, ax = plt.subplots(figsize=(10, 5.5))
    fig.patch.set_facecolor("#0a0a0a")
    ax.set_facecolor("#111827")
    bars = ax.barh(labels, values, color=["#374151", "#d4af37", "#22d3ee"])
    ax.set_xlabel("시장 규모 (USD Billion, 추정)", color="#e5e7eb")
    ax.set_title("TAM / SAM / SOM (Investor Reference)", color="#f3f4f6", fontsize=13)
    for bar, v in zip(bars, values):
        ax.text(bar.get_width() + 1, bar.get_y() + bar.get_height() / 2, f"${v}B", va="center", color="#f3f4f6")
    ax.grid(axis="x", alpha=0.25)
    p = ASSETS / "chart_market_tam.png"
    fig.tight_layout()
    fig.savefig(p, dpi=180, facecolor="#0a0a0a")
    plt.close(fig)
    return p


def chart_unit_economics() -> Path:
    metrics = ["CAC\n(₩K)", "ARPU\n(₩K/mo)", "LTV\n(6mo)", "LTV/CAC"]
    values = [40, 56.2, 337, 8.4]
    fig, ax = plt.subplots(figsize=(10, 5.5))
    fig.patch.set_facecolor("#0a0a0a")
    ax.set_facecolor("#111827")
    colors = ["#f87171", "#d4af37", "#34d399", "#22d3ee"]
    ax.bar(metrics, values, color=colors)
    ax.set_title("Unit Economics (기준 가정)", color="#f3f4f6", fontsize=13)
    ax.set_ylabel("Value", color="#e5e7eb")
    for i, v in enumerate(values):
        label = f"₩{v}K" if i < 3 else f"{v}x"
        ax.text(i, v + 5, label, ha="center", color="#f3f4f6", fontsize=11)
    ax.grid(axis="y", alpha=0.25)
    p = ASSETS / "chart_unit_economics.png"
    fig.tight_layout()
    fig.savefig(p, dpi=180, facecolor="#0a0a0a")
    plt.close(fig)
    return p


def chart_competitor_pricing() -> Path:
    """KR/US retail analysis tools — 월 환산 구독료 (2025~2026 공개 요금 기준)."""
    names = [
        "스노우볼72\n(무료)",
        "퀀터스\n(연199K÷12)",
        "TradingView\nEssential",
        "Seeking Alpha\nPremium",
        "핀트 테마\n(월정액)",
        "WallPilot\nDay",
        "WallPilot\nPremium",
        "WallPilot\nElite",
    ]
    # KRW monthly equivalent (FX $1=₩1,350)
    values = [0, 16.6, 17.5, 33.6, 4.9, 39, 99, 199]
    colors = ["#64748b", "#94a3b8", "#94a3b8", "#94a3b8", "#94a3b8", "#d4af37", "#22d3ee", "#34d399"]
    fig, ax = plt.subplots(figsize=(10, 5.8))
    fig.patch.set_facecolor("#0a0a0a")
    ax.set_facecolor("#111827")
    bars = ax.barh(names, values, color=colors)
    ax.set_xlabel("월 환산 구독료 (천 원, KRW)", color="#e5e7eb")
    ax.set_title("경쟁·유사 서비스 가격 벤치마크 (2026)", color="#f3f4f6", fontsize=13)
    for bar, v in zip(bars, values):
        if v > 0:
            ax.text(bar.get_width() + 2, bar.get_y() + bar.get_height() / 2, f"₩{v:.0f}K", va="center", color="#f3f4f6", fontsize=10)
    ax.grid(axis="x", alpha=0.25)
    ax.text(0.02, -0.12, "출처: 공개 요금(퀀터스·SA·TradingView·핀트·WallPilot) · 로보어드바이저(AUM%)는 별도", transform=ax.transAxes, color="#9ca3af", fontsize=8)
    p = ASSETS / "chart_competitor_pricing.png"
    fig.tight_layout()
    fig.savefig(p, dpi=180, facecolor="#0a0a0a", bbox_inches="tight")
    plt.close(fig)
    return p


def chart_site_structure() -> Path:
    """WallPilot Pro 사이트맵 — src/lib/membership/menus.ts 기반."""
    import matplotlib.patches as mpatches

    fig, ax = plt.subplots(figsize=(11, 6.5))
    fig.patch.set_facecolor("#0a0a0a")
    ax.set_facecolor("#0a0a0a")
    ax.set_xlim(0, 10)
    ax.set_ylim(0, 7)
    ax.axis("off")
    ax.set_title("WallPilot Pro 사이트 구조도 (Production Menu)", color="#f3f4f6", fontsize=14, pad=16)

    def box(x, y, w, h, text, color="#1f2937", ec="#d4af37"):
        rect = mpatches.FancyBboxPatch((x, y), w, h, boxstyle="round,pad=0.02", fc=color, ec=ec, lw=1.2)
        ax.add_patch(rect)
        ax.text(x + w / 2, y + h / 2, text, ha="center", va="center", color="#f3f4f6", fontsize=8, wrap=True)

    box(3.8, 6.0, 2.4, 0.55, "WallPilot Pro PWA\nwallpilotpro.vercel.app", "#111827")
    menus = [
        ("Scanner /", "Free", 0.3, 4.8),
        ("Wall St. Report", "Day+", 1.55, 4.8),
        ("DARTLAB", "Day+", 2.8, 4.8),
        ("AI Pilot", "Premium", 4.05, 4.8),
        ("Agent Desk", "Premium", 5.3, 4.8),
        ("Signal Hub", "Day+", 6.55, 4.8),
        ("RL Lab", "Elite", 7.8, 4.8),
        ("Toss Trader", "Day+", 0.3, 3.3),
        ("My API", "Free", 1.55, 3.3),
        ("Pricing", "Free", 2.8, 3.3),
    ]
    for label, tier, x, y in menus:
        box(x, y, 1.1, 0.7, f"{label}\n[{tier}]", ec="#22d3ee" if tier == "Premium" else "#d4af37" if tier == "Elite" else "#64748b")

    box(0.5, 1.5, 9.0, 1.4, "Backend: Supabase Auth/RLS · Stripe/Danal · Vercel\nData: OpenDART · 13F · Alpha Vantage · Toss API · Gemini\nIP: Terrabridge Capital · IP Shield", "#0f172a", "#374151")

    for x in [0.85, 2.1, 3.35, 4.6, 5.85, 7.1, 8.35]:
        ax.annotate("", xy=(x + 0.55, 4.75), xytext=(x + 0.55, 6.0), arrowprops=dict(arrowstyle="-", color="#374151"))
    p = ASSETS / "chart_site_structure.png"
    fig.savefig(p, dpi=180, facecolor="#0a0a0a", bbox_inches="tight")
    plt.close(fig)
    return p


def screen(name: str) -> Path:
    return SCREENS / f"{name}.png"


def add_screenshot_slide(
    prs: Presentation,
    title: str,
    img: Path,
    bullets: list[str] | None = None,
    full_bleed: bool = False,
) -> None:
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, prs)
    tb = slide.shapes.add_textbox(Inches(0.6), Inches(0.35), Inches(12), Inches(0.65))
    tb.text_frame.paragraphs[0].text = title
    tb.text_frame.paragraphs[0].font.size = Pt(24)
    tb.text_frame.paragraphs[0].font.bold = True
    tb.text_frame.paragraphs[0].font.color.rgb = GOLD
    if img.exists():
        if full_bleed:
            slide.shapes.add_picture(str(img), Inches(0.35), Inches(1.0), width=Inches(12.6))
        else:
            slide.shapes.add_picture(str(img), Inches(0.5), Inches(1.05 if bullets else 1.0), width=Inches(7.8 if bullets else 12.3))
    if bullets:
        body = slide.shapes.add_textbox(Inches(8.5), Inches(1.1), Inches(4.3), Inches(5.5))
        tf = body.text_frame
        tf.word_wrap = True
        for i, b in enumerate(bullets):
            para = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
            para.text = b
            para.font.size = Pt(12)
            para.font.color.rgb = WHITE
            para.space_after = Pt(6)
    add_footer(slide, "wallpilotpro.vercel.app · Actual product screenshot")


def set_slide_bg(slide, prs: Presentation) -> None:
    slide.background.fill.solid()
    slide.background.fill.fore_color.rgb = NAVY


def add_footer(slide, text: str = "Terrabridge Capital Inc. | Confidential") -> None:
    box = slide.shapes.add_textbox(Inches(0.5), Inches(7.0), Inches(12), Inches(0.35))
    tf = box.text_frame
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(9)
    p.font.color.rgb = GRAY


def add_title_slide(prs: Presentation, title: str, subtitle: str = "", img: Path | None = None) -> None:
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, prs)
    if img and img.exists():
        slide.shapes.add_picture(str(img), Inches(0), Inches(0), width=prs.slide_width, height=prs.slide_height)
        overlay = slide.shapes.add_shape(1, Inches(0), Inches(0), prs.slide_width, prs.slide_height)
        overlay.fill.solid()
        overlay.fill.fore_color.rgb = RGBColor(0x0A, 0x0A, 0x0A)
        overlay.line.fill.background()
        try:
            overlay.fill.transparency = 0.45
        except Exception:
            pass
    tb = slide.shapes.add_textbox(Inches(0.8), Inches(2.2), Inches(11), Inches(2))
    tf = tb.text_frame
    p = tf.paragraphs[0]
    p.text = title
    p.font.size = Pt(40)
    p.font.bold = True
    p.font.color.rgb = GOLD
    if subtitle:
        p2 = tf.add_paragraph()
        p2.text = subtitle
        p2.font.size = Pt(18)
        p2.font.color.rgb = WHITE
        p2.space_before = Pt(12)
    add_footer(slide)


def add_section(prs: Presentation, section: str) -> None:
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, prs)
    tb = slide.shapes.add_textbox(Inches(1), Inches(3), Inches(11), Inches(1.5))
    p = tb.text_frame.paragraphs[0]
    p.text = section
    p.font.size = Pt(36)
    p.font.bold = True
    p.font.color.rgb = GOLD
    p.alignment = PP_ALIGN.CENTER
    add_footer(slide)


def add_content_slide(
    prs: Presentation,
    title: str,
    bullets: list[str],
    img: Path | None = None,
    img_right: bool = True,
) -> None:
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, prs)
    tb = slide.shapes.add_textbox(Inches(0.6), Inches(0.4), Inches(12), Inches(0.8))
    p = tb.text_frame.paragraphs[0]
    p.text = title
    p.font.size = Pt(28)
    p.font.bold = True
    p.font.color.rgb = GOLD

    text_w = Inches(6.2) if img and img.exists() else Inches(12)
    text_left = Inches(0.6) if not (img and img_right) else Inches(0.6)
    if img and img.exists() and img_right:
        slide.shapes.add_picture(str(img), Inches(6.9), Inches(1.3), width=Inches(6.0))
    elif img and img.exists():
        slide.shapes.add_picture(str(img), Inches(0.6), Inches(1.3), width=Inches(5.5))
        text_left = Inches(6.4)

    body = slide.shapes.add_textbox(text_left, Inches(1.2), text_w, Inches(5.5))
    tf = body.text_frame
    tf.word_wrap = True
    for i, b in enumerate(bullets):
        para = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        para.text = b
        para.font.size = Pt(14 if len(bullets) > 6 else 16)
        para.font.color.rgb = WHITE
        para.space_after = Pt(8)
        para.level = 0
    add_footer(slide)


def add_chart_slide(prs: Presentation, title: str, chart: Path, notes: list[str] | None = None) -> None:
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, prs)
    tb = slide.shapes.add_textbox(Inches(0.6), Inches(0.35), Inches(12), Inches(0.7))
    tb.text_frame.paragraphs[0].text = title
    tb.text_frame.paragraphs[0].font.size = Pt(26)
    tb.text_frame.paragraphs[0].font.bold = True
    tb.text_frame.paragraphs[0].font.color.rgb = GOLD
    if chart.exists():
        slide.shapes.add_picture(str(chart), Inches(0.5), Inches(1.0), width=Inches(12.3))
    if notes:
        nb = slide.shapes.add_textbox(Inches(0.6), Inches(6.5), Inches(12), Inches(0.8))
        tf = nb.text_frame
        for i, n in enumerate(notes):
            para = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
            para.text = f"• {n}"
            para.font.size = Pt(11)
            para.font.color.rgb = GRAY
    add_footer(slide)


def add_table_slide(prs: Presentation, title: str, headers: list[str], rows: list[list[str]]) -> None:
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, prs)
    tb = slide.shapes.add_textbox(Inches(0.6), Inches(0.35), Inches(12), Inches(0.7))
    tb.text_frame.paragraphs[0].text = title
    tb.text_frame.paragraphs[0].font.size = Pt(26)
    tb.text_frame.paragraphs[0].font.bold = True
    tb.text_frame.paragraphs[0].font.color.rgb = GOLD
    cols, rws = len(headers), len(rows) + 1
    table = slide.shapes.add_table(rws, cols, Inches(0.6), Inches(1.2), Inches(12.1), Inches(0.4 * rws)).table
    for j, h in enumerate(headers):
        cell = table.cell(0, j)
        cell.text = h
        cell.fill.solid()
        cell.fill.fore_color.rgb = RGBColor(0x1F, 0x29, 0x37)
        for p in cell.text_frame.paragraphs:
            p.font.bold = True
            p.font.size = Pt(12)
            p.font.color.rgb = GOLD
    for i, row in enumerate(rows, start=1):
        for j, val in enumerate(row):
            cell = table.cell(i, j)
            cell.text = val
            cell.fill.solid()
            cell.fill.fore_color.rgb = RGBColor(0x11, 0x18, 0x27)
            for p in cell.text_frame.paragraphs:
                p.font.size = Pt(11)
                p.font.color.rgb = WHITE
    add_footer(slide)


def build() -> Path:
    ensure_assets()
    charts = {
        "mrr": chart_mrr_scenarios(),
        "tier": chart_tier_mix(),
        "rev": chart_revenue_streams(),
        "pnl": chart_pnl(),
        "tam": chart_market_tam(),
        "unit": chart_unit_economics(),
        "comp": chart_competitor_pricing(),
        "sitemap": chart_site_structure(),
    }
    imgs = {
        "home": screen("01-home-landing"),
        "scanner": screen("02-scanner"),
        "pricing": screen("03-pricing"),
        "report": screen("04-wall-report"),
        "dart": screen("05-dartlab"),
        "pilot": screen("06-ai-pilot"),
        "toss": screen("07-toss-trader"),
        "signals": screen("08-signals"),
        "api": screen("09-my-api"),
    }

    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)

    # 1 Cover — 실제 랜딩 Hero
    add_title_slide(
        prs,
        "WallPilot Pro™",
        "IR & Business Plan 2026 · Data-Driven Quant Analysis\nTerrabridge Capital Inc. | wallpilotpro.vercel.app",
        imgs["home"],
    )

    # 2 Disclaimer
    add_content_slide(
        prs,
        "Important Disclaimer",
        [
            "본 자료는 Terrabridge Capital Inc.의 WallPilot Pro™ 사업·투자 설명용입니다.",
            "투자자문·투자일임·집합투자업이 아니며, 증권 매매 권유가 아닙니다.",
            "경쟁사·시장 데이터는 공개 요금·리포트 기반 reference이며 추정치를 포함합니다.",
            "Confidential — 무단 복제·배포·클론 금지 (Proprietary IP).",
        ],
    )

    # 3 Executive Summary
    add_content_slide(
        prs,
        "Executive Summary",
        [
            "WallPilot Pro = Reverse-Quant + AI 주식 분석 SaaS (KR·US) — One Screen OS",
            "10개 메뉴: Scanner → Report → DART → AI Pilot → Agent → Signal → RL → Toss",
            "4-Tier B2C (₩0~₩199K/월) + B2B API · PoC · White-label",
            "90일 MRR: 보수 ₩1.7M · 기준 ₩4.6M · 낙관 ₩11.4M (ARR ₩55M 기준)",
            "규제 포지션: 참고용 분석 도구 (Not investment advice)",
        ],
        imgs["scanner"],
    )

    # 4 Problem
    add_content_slide(
        prs,
        "Problem — 개인 투자자의 정보 격차",
        [
            "차트·재무·수급·뉴스가 앱·탭·커뮤니티에 분산 → 의사결정 지연",
            "퀀터스·스노우볼은 백테스트/ETF 중심 — KR/US 개별주+수급 통합 부족",
            "로보어드바이저(핀트·토스)는 AUM % 수수료 — ‘직접 분석’ 니즈와 다름",
            "Seeking Alpha·TradingView는 US/KR·DART·토스 실행까지 One OS 아님",
        ],
    )

    # 5 Solution
    add_content_slide(
        prs,
        "Solution — WallPilot Pro (실제 제품)",
        [
            "① Reverse-Quant Scanner — 13F·퀀트·거래량 (Free preview / Day+ full)",
            "② Wall St. Report — Lynch/Greenblatt 적정가 + 수급 라벨",
            "③ DARTLAB — OpenDART 재무·공시 AI brief (KR)",
            "④ AI Pilot · Agent Desk · Signal Hub · Toss Trader · RL Lab",
        ],
        imgs["home"],
    )

    # 6 Site structure
    add_chart_slide(
        prs,
        "Product — 사이트 구조도 (wallpilotpro.vercel.app)",
        charts["sitemap"],
        ["Free: Scanner preview · Signal Hub · My API · Pricing", "Day+: Report · DART · Toss · Signals", "Premium+: AI Pilot · Agent Desk", "Elite: RL Lab · Toss Execute"],
    )

    # 7 Landing + GTM Hero
    add_screenshot_slide(
        prs,
        "Live Product — Landing & GTM Hero",
        imgs["home"],
        [
            "월가 퀀트 방식 Hero (KO/EN)",
            "3가지 가치 제안 + How it works",
            "FAQ 5개 + Terms/Privacy",
            "CTA: 무료 스캐너 미리보기",
            "Google OAuth 가입 퍼널",
        ],
    )

    # 8 Scanner
    add_screenshot_slide(
        prs,
        "Module — Reverse-Quant Scanner (/)",
        imgs["scanner"],
        [
            "KR·US 실시간 스크리닝",
            "Toss · 13F · Quant 토글",
            "Short Squeeze / High Cash 매트릭스",
            "Buying Zone · Execute (Elite)",
            "Min Tier: Free preview → Day+",
        ],
        full_bleed=not imgs["scanner"].exists(),
    )

    # 9 Wall Report
    add_screenshot_slide(
        prs,
        "Module — Wall St. Report (/wall-street-report)",
        imgs["report"],
        ["Fair Value · Buying Zone", "외국인/기관 수급 라벨", "Lynch PEG · Greenblatt ROIC", "PDF export (Premium+)", "Min Tier: Day Trading"],
    )

    # 10 DARTLAB + AI Pilot
    add_screenshot_slide(
        prs,
        "Module — DARTLAB & AI Pilot",
        imgs["dart"] if imgs["dart"].exists() else imgs["pilot"],
        [
            "DARTLAB: KR 재무·공시 AI (/dartlab)",
            "AI Pilot: Gemini 종목 Q&A (/ai-pilot)",
            "Agent Desk: Bull/Bear debate (/agents/desk)",
            "Cross-check DART 원문 권장 (면책)",
        ],
    )

    # 11 Toss + Signals
    add_screenshot_slide(
        prs,
        "Module — Toss Trader & Signal Hub",
        imgs["toss"],
        [
            "Toss Open API: 보유·손익·미체결 (/toss-trader)",
            "Signal Hub: 시그널 피드 (/signals)",
            "Elite: 분할 지정가 실행 워크플로",
            "Crypto Engine strip (로컬 연동)",
        ],
    )

    # 12 Pricing screenshot + table
    add_screenshot_slide(
        prs,
        "Monetization — Pricing (/pricing)",
        imgs["pricing"],
        [
            "Free ₩0 — Scanner preview",
            "Day ₩39K — Full scan · Report · DART",
            "Premium ₩99K — AI · Agent · PDF",
            "Elite ₩199K — Toss execute · RL Lab",
            "Stripe (US) + Danal (KR)",
        ],
    )

    # 13 Architecture
    add_table_slide(
        prs,
        "Technical Architecture (Repo Analysis)",
        ["Layer", "Stack", "Role"],
        [
            ["Frontend", "React 19 · TanStack Start · PWA · 8 langs", "Landing · Nav · Tier gates"],
            ["Engine", "Quant · Supply/Demand · Valuation · Agents", "Reverse-Quant pipeline"],
            ["Data", "OpenDART · 13F · AV · Toss · Gemini", "KR/US cross-market"],
            ["SaaS", "Supabase RLS · Stripe/Danal · Vercel", "Auth · Billing · Deploy"],
            ["Security", "IP Shield · CSP · violation log", "Clone protection"],
        ],
    )

    # 14 Competitive landscape (research-backed)
    add_table_slide(
        prs,
        "Competitive Landscape — 2026 Benchmark",
        ["Player", "Model", "Price (ref.)", "vs WallPilot"],
        [
            ["WallPilot Pro", "Flat SaaS", "₩39~199K/mo", "KR+US quant+AI+Toss OS"],
            ["퀀터스", "구독+실전권", "연 ₩199K~359K", "백테스트·자동매매 중심"],
            ["스노우볼72", "Freemium", "₩0", "ETF 자산배분 only"],
            ["Seeking Alpha", "Flat SaaS", "~$299/yr", "US research · no KR DART"],
            ["TradingView", "Flat SaaS", "$12~60/mo", "Charting · weak quant AI"],
            ["핀트/토스 AI", "AUM %", "0.55~1.17%/yr", "일임·로보 · not DIY analysis"],
            ["Bloomberg", "Terminal", "~$2K/mo", "Institutional · not retail"],
        ],
    )

    # 15 Competitor pricing chart
    add_chart_slide(
        prs,
        "Pricing Benchmark — Retail Analysis Tools (월 환산)",
        charts["comp"],
        ["WallPilot Day ₩39K = SA Premium·TradingView Pro 대비 경쟁력", "Premium ₩99K = AI Agent·PDF 포함", "로보어드바이저(AUM%)는 별도 카테고리"],
    )

    # 16 Market TAM
    add_chart_slide(
        prs,
        "Market Opportunity — TAM / SAM / SOM",
        charts["tam"],
        ["글로벌 리테일 퀀트·AI 분석 TAM ~$120B", "KR+US active trader SAM ~$18B", "WallPilot Y3 SOM target ~$2.4B"],
    )

    # 17 MRR
    add_chart_slide(
        prs,
        "Financial Projection — MRR (₩M)",
        charts["mrr"],
        ["Month 3: 보수 ₩1.7M · 기준 ₩4.6M · 낙관 ₩11.4M", "Month 12 기준: ~₩28M MRR"],
    )

    # 18 Tier mix
    add_chart_slide(prs, "Revenue Model — Tier Mix (Blended ARPU ₩56.2K)", charts["tier"])

    # 19 Unit economics
    add_chart_slide(
        prs,
        "Unit Economics",
        charts["unit"],
        ["Blended CAC ≤ ₩40K", "LTV/CAC 8.4x (6mo)", "Payback ~2 months on Day tier"],
    )

    # 20 P&L
    add_chart_slide(
        prs,
        "90-Day P&L Sketch (Base)",
        charts["pnl"],
        ["Month 3 net positive ~₩2.1M (excl. salary)", "Revenue streams: B2C + B2B PoC + US"],
    )

    # 21 KPI + GTM
    add_table_slide(
        prs,
        "90-Day KPI & GTM (Base Scenario)",
        ["Metric / Phase", "Target"],
        [
            ["Free signups", "800"],
            ["Paid (Month 3)", "80 · MRR ₩4.6M"],
            ["Free→Paid CR", "8–12%"],
            ["W1-2", "Landing · Terms · Analytics ✓"],
            ["W3-4", "SNS 3 picks · Free→Day funnel"],
            ["W9-12", "B2B 10 outreach · Elite beta"],
        ],
    )

    # 22 B2B + IP
    add_table_slide(
        prs,
        "B2B License & IP Moat",
        ["Item", "Detail"],
        [
            ["Starter API", "$2,500/mo · 1K calls/day"],
            ["Growth", "$8,000/mo · co-brand · SLA"],
            ["PoC", "$5,000 flat · 8 weeks"],
            ["IP Owner", "Terrabridge Capital Inc."],
            ["Inventor", "kangjunchul8@gmail.com"],
            ["Protection", "IP Shield · Terms · Privacy live"],
        ],
    )

    # 23 Roadmap
    add_content_slide(
        prs,
        "Roadmap — 6~18 Months",
        [
            "Q2 2026: GTM · 80+ paid · B2B PoC · wallpilotpro.vercel.app live",
            "Q3 2026: US 30% revenue · Enterprise API · GA4 funnel",
            "Q4 2026: 250+ paid · MRR ₩12M+ · Danal full",
            "2027: RIA channel · Mobile native · US trademark",
        ],
    )

    # 24 Ask
    add_content_slide(
        prs,
        "Investment / Partnership Ask",
        [
            "Strategic: KR securities · US RIA · research YouTube/Telegram",
            "PoC: 8-week embed · co-brand signal feed",
            "Optional Seed: GTM · B2B sales · compliance",
            "terrabridgecapital@gmail.com · kangjunchul8@gmail.com",
        ],
        imgs["pricing"],
    )

    # 25 Thank you
    add_title_slide(
        prs,
        "Thank You",
        "WallPilot Pro™ — wallpilotpro.vercel.app\nData-Driven Quant Analysis for Everyone",
        imgs["home"],
    )

    prs.save(str(OUT_PPTX))
    import shutil

    for f in ASSETS.glob("*"):
        if f.is_file():
            dest = OUT_DIR / f.name
            if f != dest:
                shutil.copy2(f, dest)
    if SCREENS.is_dir():
        screen_dest = OUT_DIR / "ir-screens"
        screen_dest.mkdir(parents=True, exist_ok=True)
        for f in SCREENS.glob("*.png"):
            shutil.copy2(f, screen_dest / f.name)
    return OUT_PPTX


if __name__ == "__main__":
    path = build()
    print(f"Generated: {path}")
    print(f"Assets: {OUT_DIR}")
