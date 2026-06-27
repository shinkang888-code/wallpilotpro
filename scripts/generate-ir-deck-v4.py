#!/usr/bin/env python3
"""WallPilot Pro IR Deck v4 — image-dominant, American Standard + Quant theory."""
from __future__ import annotations

import importlib.util
import shutil
from datetime import date
from pathlib import Path

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt

ROOT = Path(__file__).resolve().parents[1]
_v2 = ROOT / "scripts" / "generate-ir-deck.py"
_spec = importlib.util.spec_from_file_location("generate_ir_deck", _v2)
_mod = importlib.util.module_from_spec(_spec)
assert _spec.loader
_spec.loader.exec_module(_mod)

chart_competitor_pricing = _mod.chart_competitor_pricing
chart_market_tam = _mod.chart_market_tam
chart_mrr_scenarios = _mod.chart_mrr_scenarios
chart_pnl = _mod.chart_pnl
chart_revenue_streams = _mod.chart_revenue_streams
chart_site_structure = _mod.chart_site_structure
chart_tier_mix = _mod.chart_tier_mix
chart_unit_economics = _mod.chart_unit_economics

ASSETS = ROOT / "assets"
V4 = ASSETS / "ir-v4"
PDF = ASSETS / "ir-pdf-pages"
SCREENS = ASSETS / "ir-screens"
OUT = Path.home() / "Downloads" / f"WallPilot_Pro_IR_{date.today().isoformat()}_v4"
PPTX = OUT / "WallPilot_Pro_IR_Business_Plan_48slides_v4.pptx"

NAVY = RGBColor(0x0A, 0x0A, 0x0A)
GOLD = RGBColor(0xD4, 0xAF, 0x37)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
GRAY = RGBColor(0x9C, 0xA3, 0xAF)
CYAN = RGBColor(0x22, 0xD3, 0xEE)
SOFT = RGBColor(0x11, 0x18, 0x27)
W, H = Inches(13.333), Inches(7.5)


def img(name: str) -> Path:
    return V4 / name


def pdf(n: int) -> Path:
    return PDF / f"ir-page-{n:02d}.png"


def ensure() -> None:
    V4.mkdir(parents=True, exist_ok=True)
    OUT.mkdir(parents=True, exist_ok=True)


def bg(slide, c=NAVY) -> None:
    slide.background.fill.solid()
    slide.background.fill.fore_color.rgb = c


def foot(slide, t="Terrabridge Capital Inc. · Confidential · wallpilotpro.vercel.app") -> None:
    b = slide.shapes.add_textbox(Inches(0.4), Inches(7.06), Inches(12.5), Inches(0.3))
    b.text_frame.paragraphs[0].text = t
    b.text_frame.paragraphs[0].font.size = Pt(8)
    b.text_frame.paragraphs[0].font.color.rgb = GRAY


def line(slide, y=0.9) -> None:
    s = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.4), Inches(y), Inches(12.5), Inches(0.035))
    s.fill.solid()
    s.fill.fore_color.rgb = GOLD
    s.line.fill.background()


def ttl(slide, text: str, y=0.32, sz=28, c=GOLD, align=PP_ALIGN.LEFT) -> None:
    b = slide.shapes.add_textbox(Inches(0.5), Inches(y), Inches(12.3), Inches(0.75))
    p = b.text_frame.paragraphs[0]
    p.text = text
    p.font.size = Pt(sz)
    p.font.bold = True
    p.font.color.rgb = c
    p.alignment = align


def pic(slide, path: Path, l, t, w, h=None) -> None:
    if path.exists():
        kw = {"width": Inches(w)}
        if h:
            kw["height"] = Inches(h)
        slide.shapes.add_picture(str(path), Inches(l), Inches(t), **kw)


def overlay(slide, alpha=0.45) -> None:
    o = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), W, H)
    o.fill.solid()
    o.fill.fore_color.rgb = NAVY
    o.line.fill.background()
    try:
        o.fill.transparency = alpha
    except Exception:
        pass


def full_bleed(slide, path: Path, overlay_alpha=0.0) -> None:
    pic(slide, path, 0, 0, 13.333, 7.5)
    if overlay_alpha:
        overlay(slide, overlay_alpha)


def caption(slide, title: str, sub: str = "") -> None:
    bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0), Inches(6.05), W, Inches(1.35))
    bar.fill.solid()
    bar.fill.fore_color.rgb = NAVY
    try:
        bar.fill.transparency = 0.12
    except Exception:
        pass
    bar.line.fill.background()
    b = slide.shapes.add_textbox(Inches(0.55), Inches(6.18), Inches(12.2), Inches(1.1))
    tf = b.text_frame
    p = tf.paragraphs[0]
    p.text = title
    p.font.size = Pt(20)
    p.font.bold = True
    p.font.color.rgb = GOLD
    if sub:
        p2 = tf.add_paragraph()
        p2.text = sub
        p2.font.size = Pt(11)
        p2.font.color.rgb = GRAY
        p2.space_before = Pt(3)


def section(prs, path: Path, title: str, sub: str) -> None:
    s = prs.slides.add_slide(prs.slide_layouts[6])
    bg(s)
    full_bleed(s, path, 0.38)
    ttl(s, title, y=2.65, sz=42, align=PP_ALIGN.CENTER)
    b = s.shapes.add_textbox(Inches(1.2), Inches(3.55), Inches(10.9), Inches(0.7))
    b.text_frame.paragraphs[0].text = sub
    b.text_frame.paragraphs[0].font.size = Pt(15)
    b.text_frame.paragraphs[0].font.color.rgb = CYAN
    b.text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
    line(s, 3.35)
    foot(s)


def theory_split(slide, path: Path, title: str, paras: list[str], img_w=7.2) -> None:
    bg(slide)
    ttl(slide, title, sz=24)
    line(slide)
    pic(slide, path, 0.4, 1.0, img_w, 6.0)
    b = slide.shapes.add_textbox(Inches(7.85), Inches(1.05), Inches(5.1), Inches(5.9))
    tf = b.text_frame
    tf.word_wrap = True
    for i, t in enumerate(paras):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.text = t
        p.font.size = Pt(11)
        p.font.color.rgb = WHITE
        p.space_after = Pt(10)
        p.line_spacing = 1.12
    foot(slide)


def kpi_row(slide, items: list[tuple[str, str, str]], y=1.15) -> None:
    gap = 12.3 / len(items)
    for i, (lb, val, note) in enumerate(items):
        x = 0.45 + i * gap
        c = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(x), Inches(y), Inches(gap - 0.12), Inches(1.45))
        c.fill.solid()
        c.fill.fore_color.rgb = SOFT
        c.line.color.rgb = GOLD
        c.line.width = Pt(0.6)
        b = slide.shapes.add_textbox(Inches(x + 0.1), Inches(y + 0.12), Inches(gap - 0.28), Inches(1.2))
        tf = b.text_frame
        p = tf.paragraphs[0]
        p.text = lb
        p.font.size = Pt(9)
        p.font.color.rgb = GRAY
        p2 = tf.add_paragraph()
        p2.text = val
        p2.font.size = Pt(22)
        p2.font.bold = True
        p2.font.color.rgb = GOLD
        p3 = tf.add_paragraph()
        p3.text = note
        p3.font.size = Pt(8)
        p3.font.color.rgb = CYAN


def table_sl(slide, title: str, headers: list[str], rows: list[list[str]]) -> None:
    bg(slide)
    ttl(slide, title, sz=22)
    line(slide)
    tbl = slide.shapes.add_table(len(rows) + 1, len(headers), Inches(0.45), Inches(1.05), Inches(12.4), Inches(min(0.36 * (len(rows) + 1), 5.5))).table
    for j, h in enumerate(headers):
        c = tbl.cell(0, j)
        c.text = h
        c.fill.solid()
        c.fill.fore_color.rgb = RGBColor(0x1F, 0x29, 0x37)
        for p in c.text_frame.paragraphs:
            p.font.bold = True
            p.font.size = Pt(10)
            p.font.color.rgb = GOLD
    for i, row in enumerate(rows, 1):
        for j, v in enumerate(row):
            c = tbl.cell(i, j)
            c.text = v
            c.fill.solid()
            c.fill.fore_color.rgb = SOFT if i % 2 else RGBColor(0x0F, 0x17, 0x2A)
            for p in c.text_frame.paragraphs:
                p.font.size = Pt(9)
                p.font.color.rgb = WHITE
    foot(slide)


def chart_sl(slide, title: str, chart: Path, notes: list[str] | None = None) -> None:
    bg(slide)
    ttl(slide, title, sz=22)
    pic(slide, chart, 0.3, 0.92, 12.75)
    if notes:
        b = slide.shapes.add_textbox(Inches(0.5), Inches(6.5), Inches(12.3), Inches(0.7))
        tf = b.text_frame
        for i, n in enumerate(notes):
            p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
            p.text = f"▸ {n}"
            p.font.size = Pt(9)
            p.font.color.rgb = GRAY
    foot(slide)


def hero_product(slide, path: Path, title: str, sub: str) -> None:
    bg(slide, RGBColor(0xF1, 0xF5, 0xF9))
    frame = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.25), Inches(0.45), Inches(12.85), Inches(6.55))
    frame.fill.solid()
    frame.fill.fore_color.rgb = WHITE
    frame.line.color.rgb = RGBColor(0xE2, 0xE8, 0xF0)
    pic(slide, path, 0.35, 0.55, 12.65, 5.85)
    caption(slide, title, sub)
    foot(slide, "WallPilot Pro™ · Product UI (AI-rendered + production reference)")


def build() -> Path:
    ensure()
    charts = {k: fn() for k, fn in [
        ("mrr", chart_mrr_scenarios), ("tier", chart_tier_mix), ("rev", chart_revenue_streams),
        ("pnl", chart_pnl), ("tam", chart_market_tam), ("unit", chart_unit_economics),
        ("comp", chart_competitor_pricing), ("map", chart_site_structure),
    ]}
    prs = Presentation()
    prs.slide_width, prs.slide_height = W, H

    # 1 Cover
    s = prs.slides.add_slide(prs.slide_layouts[6])
    bg(s)
    full_bleed(s, img("ir-v4-cover-wallstreet.png"), 0.48)
    b = s.shapes.add_textbox(Inches(0.65), Inches(1.35), Inches(11), Inches(2.2))
    tf = b.text_frame
    tf.paragraphs[0].text = "WallPilot Pro™"
    tf.paragraphs[0].font.size = Pt(54)
    tf.paragraphs[0].font.bold = True
    tf.paragraphs[0].font.color.rgb = GOLD
    p = tf.add_paragraph()
    p.text = "American Standard · Reverse-Quant · KR × US"
    p.font.size = Pt(20)
    p.font.color.rgb = WHITE
    p.space_before = Pt(8)
    kpi_row(s, [
        ("M3 MRR", "₩4.6M", "Base · 80 paid"),
        ("ARR", "₩55M", "Run-rate"),
        ("LTV/CAC", "8.4×", "Unit econ"),
        ("Tier", "₩0–199K", "4-Tier SaaS"),
    ], y=4.85)
    foot(s, f"IR v4 · {date.today().isoformat()} · Terrabridge Capital Inc.")

    # 2 Disclaimer
    s = prs.slides.add_slide(prs.slide_layouts[6])
    bg(s)
    ttl(s, "Important Disclaimer", sz=26)
    line(s)
    b = s.shapes.add_textbox(Inches(0.55), Inches(1.15), Inches(12.2), Inches(5.5))
    tf = b.text_frame
    for i, t in enumerate([
        "Terrabridge Capital Inc. · WallPilot Pro™ IR & Partnership Materials",
        "Not investment advice · Not RIA · Not broker-dealer · Reference analysis only",
        "Financial projections = GTM scenarios · Public pricing benchmarks · Not audited",
        "Proprietary IP · IP Shield · Unauthorized clone/scrape/embed prohibited",
    ]):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.text = f"• {t}"
        p.font.size = Pt(14)
        p.font.color.rgb = WHITE
        p.space_after = Pt(12)
    foot(s)

    # 3–11 THEORY BLOCK
    section(prs, img("ir-v4-american-standard.png"),
            "00 · American Standard", "Institutional-grade quant logic for global retail")

    s = prs.slides.add_slide(prs.slide_layouts[6])
    full_bleed(s, img("ir-v4-american-standard.png"), 0.35)
    caption(s, "American Standard — What We Export to Retail",
            "SEC/GAAP disclosure culture · 13F transparency · Factor models · Risk gates · Not gambling — process")

    s = prs.slides.add_slide(prs.slide_layouts[6])
    full_bleed(s, img("ir-v4-symbol-wallstreet.png"), 0.25)
    caption(s, "Wall Street Symbolism → Product Philosophy",
            "Bull = risk-on momentum · Exchange = price discovery · 13F = smart-money footprint · Ticker = real-time truth")

    section(prs, img("ir-v4-reverse-quant-theory.png"),
            "01 · Reverse-Quant Theory", "Follow the giants before the quote board catches up")

    s = prs.slides.add_slide(prs.slide_layouts[6])
    theory_split(s, img("ir-v4-reverse-quant-theory.png"), "Reverse-Quant Engine", [
        "Forward prediction fails retail investors. Reverse-Quant inverts the problem: "
        "start from what elite capital already did — 13F whales, Toss high-rollers, short interest.",
        "Lag window: institutional footprint → public price. WallPilot compresses 13F + volume + "
        "supply/demand into one screener matrix (Short Squeeze · Dividend Compounders).",
        "KR × US: same engine, dual universe (KOSPI/KOSDAQ + NYSE/NASDAQ). "
        "One Screen OS — no tab switching across chart, filing, and broker apps.",
        "Output: OVERWEIGHT/HOLD/UNDERWEIGHT + Magic D quant grade + 30D momentum — "
        "reference zones, not trade orders.",
    ])

    s = prs.slides.add_slide(prs.slide_layouts[6])
    theory_split(s, img("ir-v4-quant-strategy.png"), "Quantum Investment Strategy (WallPilot Framework)", [
        "Metaphor: markets as probabilistic systems — multiple factors, overlapping scenarios, "
        "not single-indicator bets. WallPilot applies multi-pillar scoring (4 Pillars in AI Pilot).",
        "Factor stack: (1) Neglected bottom · (2) Cash fortress · (3) Megatrend alignment · "
        "(4) Catalyst within 4 weeks — filters noise before narrative.",
        "Valuation layer: Peter Lynch GARP (PEG, stalwart/growth type) + Joel Greenblatt "
        "Magic Formula (ROIC × Earnings Yield) → fair value & margin of safety.",
        "Risk gate: Bull/Bear agent debate → Verdict → Hard stop & split-limit execution guide. "
        "Past model output ≠ future returns — disclosed on every report.",
    ])

    s = prs.slides.add_slide(prs.slide_layouts[6])
    theory_split(s, img("ir-v4-ui-wall-report.png"), "Classic Quant Models — Built In", [
        "Peter Lynch: PEG ratio, Lynch Score, growth type (fast-grower / stalwart / turnaround).",
        "Joel Greenblatt: ROIC, Earnings Yield, Magic Grade — quality × cheapness ranking.",
        "Supply/Demand (KR): foreign & institutional net flow labels when MCP/broker API connected.",
        "Technical overlay: RSI · MACD · Bollinger %B on every Wall St. Report card.",
        "Agent Deep Report: Research Manager persona · KR execution guide · catalyst calendar.",
    ])

    s = prs.slides.add_slide(prs.slide_layouts[6])
    bg(s)
    ttl(s, "Smart Money Data Stack", sz=24)
    line(s)
    kpi_row(s, [
        ("13F", "Whales", "Quarterly institutional"),
        ("Toss HR", "Top 20", "KR high-roller flow"),
        ("Short", "Squeeze", "Interest vs volume"),
        ("DART", "KR Filings", "OpenDART AI brief"),
    ])
    pic(s, pdf(1), 0.45, 2.85, 12.4, 4.0)
    foot(s)

    # 12 Vision
    section(prs, img("ir-v4-section-product.png"), "02 · Vision & Opportunity", "One Screen OS · B2C + B2B · KR + US")

    s = prs.slides.add_slide(prs.slide_layouts[6])
    bg(s)
    ttl(s, "Executive Summary", sz=26)
    line(s)
    kpi_row(s, [("Product", "10 Menus", "Scan→Execute"), ("M3 MRR", "₩4.6M", "80 paid"), ("B2B", "$2.5K+", "API license"), ("Reg", "Not IA", "Ref tool")])
    pic(s, img("ir-v4-ui-scanner.png"), 0.4, 2.75, 12.5, 4.15)
    foot(s)

    s = prs.slides.add_slide(prs.slide_layouts[6])
    bg(s)
    ttl(s, "North Star · Unit Economics", sz=24)
    line(s)
    kpi_row(s, [
        ("KPI", "≥2 Scan/wk", "Paid WAU"),
        ("CR", "8–12%", "Free→Paid"),
        ("ARPU", "₩56.2K", "Blended"),
        ("Payback", "~2 mo", "Day tier"),
    ])
    b = s.shapes.add_textbox(Inches(0.55), Inches(2.85), Inches(12.2), Inches(3.8))
    tf = b.text_frame
    for i, t in enumerate([
        "CAC blended ≤ ₩40K · LTV(6mo) ₩337K · LTV/CAC 8.4× · Churn 10%/mo (base)",
        "Google OAuth → Free/active · Stripe/Danal → pro/premium/elite",
        "Admin: shinkang888 · Product: kangjunchul8 · Business: terrabridgecapital@gmail.com",
    ]):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.text = f"• {t}"
        p.font.size = Pt(13)
        p.font.color.rgb = WHITE
        p.space_after = Pt(10)
    foot(s)

    s = prs.slides.add_slide(prs.slide_layouts[6])
    theory_split(s, img("ir-v4-ui-scanner.png"), "Problem — Fragmented Retail Stack", [
        "4+ apps for chart, news, filings, broker — decision latency kills edge.",
        "Quantors/Snowball: backtest & ETF — weak KR/US cross + supply integration.",
        "Robo-advisors (Pint/Toss AI): AUM fee model ≠ DIY analyst workflow.",
        "Seeking Alpha / TradingView: no DART · no Toss execute · not One OS.",
    ], img_w=6.8)

    s = prs.slides.add_slide(prs.slide_layouts[6])
    table_sl(s, "Pain Points vs WallPilot", ["Pain", "Impact", "WallPilot"], [
        ["Info silos", "4+ app switches", "One Screen OS"],
        ["KR·US cross", "Dual tracking", "Unified universe"],
        ["13F lag", "Raw data only", "Reverse-Quant scan"],
        ["AI trust", "Hallucination", "DART cross-check + Risk Gate"],
        ["Execution gap", "Analysis→order friction", "Toss Trader Elite"],
    ])

    s = prs.slides.add_slide(prs.slide_layouts[6])
    bg(s)
    ttl(s, "Solution Pipeline", sz=26)
    line(s)
    kpi_row(s, [
        ("① Scan", "Reverse-Quant", "13F·Vol·Quant"),
        ("② Report", "Lynch·Greenblatt", "Fair value"),
        ("③ AI", "AlphaQuant", "Bull/Bear"),
        ("④ Execute", "Toss API", "Elite"),
    ])
    pic(s, charts["map"], 0.35, 2.85, 12.65, 3.95)
    foot(s)

    # 19 Product section
    section(prs, img("ir-v4-section-product.png"), "03 · Live Product", "AI-rendered UI + production reference")

    product_slides = [
        (img("ir-v4-ui-scanner.png"), "Module — Reverse-Quant Scanner", "Live Signals · Short Squeeze · OVERWEIGHT/HOLD · Magic D · 30D Momentum · KR+US"),
        (img("ir-v4-ui-scanner.png"), "Scanner — Guru Tracker Hero", "Top 20 Toss high-rollers · 13F whales · Super quant insiders · NYSE/NASDAQ · KOSPI/KOSDAQ"),
        (pdf(2) if pdf(2).exists() else img("ir-v4-ui-scanner.png"), "Production — Live Signals Matrix", "wallpilotir.pdf reference · 20+ tickers · dual screener columns"),
        (img("ir-v4-ui-wall-report.png"), "Module — Wall St. Report", "Fair Value · Buying Zone · Lynch PEG · Greenblatt ROIC · RSI/MACD/BB · 3M chart"),
        (pdf(6) if pdf(6).exists() else img("ir-v4-ui-wall-report.png"), "Agent Deep Report", "Bull vs Bear · Risk Gate · KR execution guide · Catalyst calendar"),
        (img("ir-v4-ui-ai-pilot.png"), "Module — AI Pilot", "AlphaQuant chat · 4 Pillars · Scanner-linked context · Premium+"),
        (pdf(4) if pdf(4).exists() else img("ir-v4-ui-ai-pilot.png"), "AI Trade Setup Cards", "Buy range · Take profit · Hard stop · 20 picks · RSI/MACD"),
        (img("ir-v4-ui-dartlab.png"), "Module — DARTLAB", "OpenDART filings · AI brief · PDF export · Day+ tier"),
        (img("ir-v4-ui-toss-trader.png"), "Module — Toss Trader", "Holdings · PnL · Buying power · Open orders · Elite execute"),
        (img("ir-v4-ui-signals.png"), "Module — Signal Hub", "Community signals · copy workflow · Free read · Day+ write"),
        (SCREENS / "01-home-landing.png" if (SCREENS / "01-home-landing.png").exists() else img("ir-v4-ui-scanner.png"),
         "Landing & GTM Hero", "3 Value Props · FAQ · Terms/Privacy · CTA · Vercel Analytics"),
        (SCREENS / "09-my-api.png" if (SCREENS / "09-my-api.png").exists() else img("ir-v4-ui-pricing.png"),
         "Module — My API", "Toss Open API key · connection test · encrypted · Free tier"),
        (img("ir-v4-ui-pricing.png"), "Monetization — Pricing", "Free ₩0 · Day ₩39K · Premium ₩99K · Elite ₩199K · Stripe + Danal"),
    ]
    for path, title, sub in product_slides:
        s = prs.slides.add_slide(prs.slide_layouts[6])
        use_hero = "ui-" in path.name or path.parent.name == "ir-v4"
        if use_hero and "ir-pdf" not in str(path):
            hero_product(s, path, title, sub)
        else:
            full_bleed(s, path, 0.08)
            caption(s, title, sub)
            foot(s)

    s = prs.slides.add_slide(prs.slide_layouts[6])
    table_sl(s, "Entitlement Matrix", ["Menu", "Free", "Day", "Premium", "Elite"], [
        ["Scanner", "preview", "execute", "execute+PDF", "execute+PDF"],
        ["Wall Report", "—", "execute", "execute+PDF", "execute+PDF"],
        ["DARTLAB", "view", "execute+PDF", "execute+PDF", "execute+PDF"],
        ["AI Pilot", "—", "—", "execute+PDF", "execute+PDF"],
        ["Toss Trader", "view", "dashboard", "dashboard", "orders"],
        ["RL Lab", "—", "—", "—", "execute+PDF"],
    ])

    # Business
    section(prs, img("ir-v4-section-finance.png"), "04 · Business Model", "4-Tier B2C · B2B API · Competitive moat")

    s = prs.slides.add_slide(prs.slide_layouts[6])
    full_bleed(s, img("ir-v4-ui-pricing.png"), 0.05)
    caption(s, "Revenue Architecture", "B2C subscription · B2B API · PoC · US Stripe · Blended ARPU ₩56.2K")
    foot(s)

    s = prs.slides.add_slide(prs.slide_layouts[6])
    table_sl(s, "Competitive Landscape 2026", ["Player", "Model", "Price", "Gap"], [
        ["WallPilot", "Flat SaaS", "₩39~199K", "KR+US OS"],
        ["Quantors", "Sub+live", "₩199K~359K/yr", "Auto-trade focus"],
        ["Snowball72", "Free", "₩0", "ETF only"],
        ["Seeking Alpha", "Flat", "~$299/yr", "US only"],
        ["TradingView", "Flat", "$12~60/mo", "Charts"],
        ["Bloomberg", "Terminal", "~$2K/mo", "Institutional"],
    ])

    for title, key, notes in [
        ("Pricing Benchmark (₩K/mo)", "comp", ["Day ₩39K competitive vs SA/TV", "Premium bundles AI+PDF"]),
        ("TAM / SAM / SOM", "tam", ["TAM $120B · SAM $18B · SOM $2.4B Y3"]),
        ("MRR Scenarios (₩M)", "mrr", ["M3 base ₩4.6M · M12 base ₩28M"]),
        ("Tier Mix · ARPU", "tier", ["Day 60% · Premium 30% · Elite 10%"]),
        ("Revenue Streams", "rev", ["M3: B2C+PoC+US · M12: B2C+B2B"]),
        ("Unit Economics", "unit", ["LTV/CAC 8.4× · Payback 2mo"]),
        ("90-Day P&L (Base)", "pnl", ["M3 net +₩2.1M excl salary"]),
    ]:
        s = prs.slides.add_slide(prs.slide_layouts[6])
        chart_sl(s, title, charts[key], notes)

    s = prs.slides.add_slide(prs.slide_layouts[6])
    table_sl(s, "B2B License", ["Tier", "USD/mo", "Includes"], [
        ["Starter", "$2,500", "1K calls/day"],
        ["Growth", "$8,000", "10K/day · SLA"],
        ["Enterprise", "$25,000+", "White-label"],
        ["PoC 8wk", "$5,000", "Sandbox embed"],
    ])

    s = prs.slides.add_slide(prs.slide_layouts[6])
    full_bleed(s, img("ir-v4-reverse-quant-theory.png"), 0.2)
    caption(s, "B2B API — Embed Reverse-Quant in Partner Channels",
            "MTS · YouTube research · RIA · Fintech · Starter $2.5K · Growth $8K · PoC $5K")
    foot(s)

    s = prs.slides.add_slide(prs.slide_layouts[6])
    table_sl(s, "90-Day GTM KPI", ["Metric", "Target"], [
        ["Free signups", "800"], ["Paid M3", "80 · MRR ₩4.6M"],
        ["Free→Paid CR", "8–12%"], ["B2B PoC", "1"],
        ["SNS", "3 posts/week"], ["CAC", "≤ ₩40K"],
    ])

    s = prs.slides.add_slide(prs.slide_layouts[6])
    bg(s)
    ttl(s, "GTM Phases W1–12", sz=24)
    line(s)
    kpi_row(s, [
        ("P1 W1-2", "Landing", "Terms·Analytics"),
        ("P2 W3-4", "SNS", "Free→Day"),
        ("P3 W5-8", "Trust", "AI demo·cases"),
        ("P4 W9-12", "B2B·US", "PoC·Elite beta"),
    ])
    pic(s, img("ir-v4-section-finance.png"), 0.4, 2.85, 12.5, 3.95)
    foot(s)

    s = prs.slides.add_slide(prs.slide_layouts[6])
    table_sl(s, "Tech Stack", ["Layer", "Stack", "Role"], [
        ["Frontend", "React19·TanStack·PWA·8lang", "UI·gates"],
        ["Engine", "Quant·Agents·Valuation", "Pipeline"],
        ["Data", "DART·13F·Toss·Gemini", "KR·US"],
        ["SaaS", "Supabase·Stripe·Vercel", "Auth·bill"],
        ["Security", "IP Shield·CSP", "Anti-clone"],
    ])

    s = prs.slides.add_slide(prs.slide_layouts[6])
    theory_split(s, img("ir-v4-cover-wallstreet.png"), "Roadmap · IP · Partnership Ask", [
        "Q2 2026: 80+ paid · B2B PoC · wallpilotpro.vercel.app",
        "Q3 2026: US 30% revenue · Enterprise API",
        "Q4 2026: 250+ paid · MRR ₩12M+",
        "2027: RIA channel · Mobile · US trademark",
        "IP: Terrabridge Capital · WallPilot Pro™ · IP Shield enforced",
        "Ask: Strategic partner · 8-week PoC · Optional seed for GTM/B2B",
        "Contact: terrabridgecapital@gmail.com · kangjunchul8@gmail.com",
    ], img_w=6.5)

    s = prs.slides.add_slide(prs.slide_layouts[6])
    bg(s)
    full_bleed(s, img("ir-v4-cover-wallstreet.png"), 0.52)
    b = s.shapes.add_textbox(Inches(0.7), Inches(2.4), Inches(11.5), Inches(2.5))
    tf = b.text_frame
    tf.paragraphs[0].text = "Thank You"
    tf.paragraphs[0].font.size = Pt(50)
    tf.paragraphs[0].font.bold = True
    tf.paragraphs[0].font.color.rgb = GOLD
    p = tf.add_paragraph()
    p.text = "WallPilot Pro™ — American Standard Quant OS for Everyone"
    p.font.size = Pt(18)
    p.font.color.rgb = WHITE
    p.space_before = Pt(12)
    foot(s)

    prs.save(str(PPTX))
    shutil.copytree(V4, OUT / "ir-v4", dirs_exist_ok=True)
    for f in ASSETS.glob("chart_*.png"):
        shutil.copy2(f, OUT / f.name)
    return PPTX


if __name__ == "__main__":
    p = build()
    from pptx import Presentation as P
    print(f"Generated: {p}")
    print(f"Slides: {len(P(str(p)).slides)}")
    print(f"Assets: {OUT}")
