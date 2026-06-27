#!/usr/bin/env python3
"""Generate WallPilot Pro 15-page User Manual (PPTX) using wallpilotir.pdf screenshots."""
from __future__ import annotations

from datetime import date
from pathlib import Path

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt

ROOT = Path(__file__).resolve().parents[1]
ASSETS = ROOT / "assets"
PDF_PAGES = ASSETS / "ir-pdf-pages"
SCREENS = ASSETS / "ir-screens"
DOWNLOADS = Path.home() / "Downloads"
OUT_DIR = DOWNLOADS / f"WallPilot_Pro_User_Manual_{date.today().isoformat()}"
OUT_PPTX = OUT_DIR / "WallPilot_Pro_User_Manual_15pages.pptx"
OUT_MD = OUT_DIR / "WallPilot_Pro_User_Manual_15pages.md"

NAVY = RGBColor(0x0A, 0x0A, 0x0A)
GOLD = RGBColor(0xD4, 0xAF, 0x37)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
GRAY = RGBColor(0x9C, 0xA3, 0xAF)
CYAN = RGBColor(0x22, 0xD3, 0xEE)


def set_slide_bg(slide) -> None:
    slide.background.fill.solid()
    slide.background.fill.fore_color.rgb = NAVY


def add_footer(slide, text: str = "WallPilot Pro™ · Terrabridge Capital Inc. · 참고용 분석 도구 (투자자문 아님)") -> None:
    box = slide.shapes.add_textbox(Inches(0.5), Inches(7.05), Inches(12.3), Inches(0.35))
    p = box.text_frame.paragraphs[0]
    p.text = text
    p.font.size = Pt(8)
    p.font.color.rgb = GRAY


def add_title_slide(prs: Presentation, title: str, subtitle: str, meta: str) -> None:
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide)
    hero = PDF_PAGES / "ir-page-01.png"
    if hero.exists():
        slide.shapes.add_picture(str(hero), Inches(0), Inches(0), width=prs.slide_width, height=prs.slide_height)
        overlay = slide.shapes.add_shape(1, Inches(0), Inches(0), prs.slide_width, prs.slide_height)
        overlay.fill.solid()
        overlay.fill.fore_color.rgb = RGBColor(0x0A, 0x0A, 0x0A)
        overlay.line.fill.background()
        try:
            overlay.fill.transparency = 0.55
        except Exception:
            pass
    tb = slide.shapes.add_textbox(Inches(0.8), Inches(2.0), Inches(11.5), Inches(3))
    tf = tb.text_frame
    p = tf.paragraphs[0]
    p.text = title
    p.font.size = Pt(38)
    p.font.bold = True
    p.font.color.rgb = GOLD
    p2 = tf.add_paragraph()
    p2.text = subtitle
    p2.font.size = Pt(20)
    p2.font.color.rgb = WHITE
    p2.space_before = Pt(14)
    p3 = tf.add_paragraph()
    p3.text = meta
    p3.font.size = Pt(13)
    p3.font.color.rgb = CYAN
    p3.space_before = Pt(20)
    add_footer(slide)


def add_manual_slide(
    prs: Presentation,
    page: int,
    title: str,
    paragraphs: list[str],
    img: Path | None = None,
    bullets: list[str] | None = None,
) -> None:
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide)

    header = slide.shapes.add_textbox(Inches(0.55), Inches(0.3), Inches(12), Inches(0.75))
    hp = header.text_frame.paragraphs[0]
    hp.text = f"{page}. {title}"
    hp.font.size = Pt(24)
    hp.font.bold = True
    hp.font.color.rgb = GOLD

    has_img = img is not None and img.exists()
    text_left = Inches(0.55)
    text_w = Inches(6.0) if has_img else Inches(12.2)
    body_top = Inches(1.05)

    if has_img:
        slide.shapes.add_picture(str(img), Inches(6.75), Inches(0.95), width=Inches(6.1))

    box = slide.shapes.add_textbox(text_left, body_top, text_w, Inches(5.85))
    tf = box.text_frame
    tf.word_wrap = True

    idx = 0
    for para_text in paragraphs:
        para = tf.paragraphs[0] if idx == 0 else tf.add_paragraph()
        para.text = para_text
        para.font.size = Pt(11)
        para.font.color.rgb = WHITE
        para.space_after = Pt(10)
        para.line_spacing = 1.15
        idx += 1

    if bullets:
        for b in bullets:
            para = tf.add_paragraph()
            para.text = f"• {b}"
            para.font.size = Pt(10.5)
            para.font.color.rgb = GRAY
            para.space_after = Pt(5)
            para.level = 0

    add_footer(slide, f"페이지 {page}/15 · wallpilotpro.vercel.app")


def build_markdown(pages: list[dict]) -> str:
    lines = [
        "# WallPilot Pro 이용자 매뉴얼 (15페이지)",
        "",
        f"버전: {date.today().isoformat()} · Terrabridge Capital Inc.",
        "",
        "> 본 문서는 투자 참고용 정보·분석 도구 이용 안내입니다. 투자자문·투자일임·집합투자업에 해당하지 않습니다.",
        "",
    ]
    for p in pages:
        lines.append(f"## {p['page']}. {p['title']}")
        lines.append("")
        for para in p["paragraphs"]:
            lines.append(para)
            lines.append("")
        if p.get("bullets"):
            for b in p["bullets"]:
                lines.append(f"- {b}")
            lines.append("")
        lines.append("---")
        lines.append("")
    return "\n".join(lines)


def main() -> None:
    OUT_DIR.mkdir(parents=True, exist_ok=True)

    pages: list[dict] = [
        {
            "page": 2,
            "title": "WallPilot Pro란?",
            "img": None,
            "paragraphs": [
                "WallPilot Pro는 한국(KOSPI·KOSDAQ)과 미국(NYSE·NASDAQ) 상장 종목을 대상으로, "
                "기관 13F 공시·토스 하이롤러 흐름·공매도 잔고·거래량·재무제표·뉴스를 데이터로 종합해 "
                "“지금 주목할 종목”과 적정가·매수 관심 구간·리스크 시나리오를 참고 정보로 제공하는 "
                "Reverse-Quant 분석 플랫폼입니다.",
                "서비스 URL: https://wallpilotpro.vercel.app · PWA 형태로 PC·모바일 브라우저에서 이용합니다.",
                "핵심 가치는 “거인들이 사는 것을, 시세표가 따라잡기 전에” 데이터로 압축해 한 화면에서 "
                "판단할 수 있게 돕는 것입니다. AI는 보조 해석 도구이며, 최종 투자 판단과 손실 책임은 "
                "전적으로 이용자에게 있습니다.",
            ],
            "bullets": [
                "Reverse-Quant 스캔: 13F·퀀트·거래량 필터",
                "월가리포트: Lynch·Greenblatt식 가치평가 + AI 해설",
                "AI Pilot: 대화형 종목 분석·촉매·진입/손절 가이드",
                "토스 API 연동(Elite): 계좌 확인·주문 워크플로",
            ],
        },
        {
            "page": 3,
            "title": "회원가입·로그인·구독",
            "img": SCREENS / "03-pricing.png",
            "paragraphs": [
                "① Google 계정으로 로그인하면 Free 등급이 자동 부여됩니다. "
                "무료 회원은 스캐너 미리보기와 시그널 허브 열람이 가능합니다.",
                "② 유료 기능은 Stripe(해외 카드) 또는 Danal(국내 결제) 월 구독으로 승급합니다. "
                "결제 실패·구독 취소 시 Free(inactive)로 자동 다운그레이드됩니다.",
                "③ 상단 네비게이션의 「요금제」에서 플랜을 비교하고 구독을 시작할 수 있습니다. "
                "각 메뉴는 최소 등급이 정해져 있으며, 부족할 경우 업그레이드 안내가 표시됩니다.",
            ],
            "bullets": [
                "Free: ₩0 — 스캐너 미리보기, 시그널 허브 열람",
                "Day Trading: 월 ₩39,000 — 전체 스캔, 월가리포트, DARTLAB, 시그널 허브",
                "Premium: 월 ₩99,000 — AI Pilot, Agent Desk, PDF 출력",
                "Elite: 월 ₩199,000 — RL Lab, 토스 주문 실행 포함",
            ],
        },
        {
            "page": 4,
            "title": "화면 구조 및 메뉴",
            "img": SCREENS / "01-home-landing.png",
            "paragraphs": [
                "상단 바: WallPilot 로고, 주요 메뉴, 총 자산(토스 연동 시), Toss 바로가기, EN/KO 언어 전환.",
                "좌측 또는 상단 「토스 API 미연결」 표시는 즐권 API 키 등록 여부를 나타냅니다. "
                "My API 메뉴에서 토스 Open API를 연결하면 실시간 잔고·주문 기능을 사용할 수 있습니다.",
                "메뉴별 최소 등급: Scanner·My API·Pricing(Free) → Wall St. Report·DARTLAB·Signal Hub·Toss Trader(Day+) "
                "→ AI Pilot·Agent Desk(Premium) → RL Lab(Elite).",
            ],
            "bullets": [
                "스캐너(/) — Reverse-Quant Live Signals",
                "Wall St. Report — 종목별 퀀트 리포트",
                "DARTLAB — 공시·재무 DART 연동",
                "AI Pilot — AlphaQuant 대화형 분석",
                "Agent Desk / Signal Hub / RL Lab / Toss Trader",
            ],
        },
        {
            "page": 5,
            "title": "리버스 퀀트 엔진 개요",
            "img": PDF_PAGES / "ir-page-01.png",
            "paragraphs": [
                "스캐너 첫 화면은 「리버스 퀀트 엔진」 히어로 섹션으로 시작합니다. "
                "REVERSE-QUANT · KR · US 태그는 한·미 교차 분석을 의미합니다.",
                "「상위 20 구루 포트폴리오 트래커」 카드: 토스 상위 20 하이롤러, 월스트리트 13F 고래, "
                "슈퍼 퀀트 인사이더 세 축으로 엘리트 운용자·개인 투자자의 포지션 변화를 추적합니다.",
                "「월스트리트의 엣지」 카드: 기관 13F, 토스 하이롤러, 공매도 잔고를 하나의 결정으로 응축. "
                "하단 미국(NYSE·NASDAQ)·한국(KOSPI·KOSDAQ) 시장 카드를 눌러 스캔 유니버스를 전환합니다.",
                "「리버스 퀀트 엔진 스캔」 버튼을 누르면 Live Signals 목록으로 스크롤됩니다.",
            ],
        },
        {
            "page": 6,
            "title": "스캐너 — Live Signals",
            "img": PDF_PAGES / "ir-page-02.png",
            "paragraphs": [
                "Live Signals 영역은 실시간에 가깝게 갱신되는 종목 후보 목록입니다. "
                "좌·우 컬럼으로 스크리너 유형이 나뉩니다(예: 숏 스퀴즈 타겟, 고배당 컴파운더).",
                "각 종목 카드에는 티커·시장(KR/US), 투자의견 배지(OVERWEIGHT/HOLD/UNDERWEIGHT), "
                "Magic D 퀀트 신호, 회사명, 현재가, 30일 모멘텀(±%)이 표시됩니다.",
                "카드를 클릭하면 상세 분석·차트·월가리포트로 이어지는 워크플로를 진행할 수 있습니다. "
                "데이터 출처: Toss Open API, Yahoo Finance(미국), 자체 퀀트 엔진.",
            ],
            "bullets": [
                "OVERWEIGHT(파랑): 상대적 비중 확대 관점",
                "HOLD(노랑): 중립·관망",
                "UNDERWEIGHT(주황): 비중 축소·주의",
                "30D Momentum: 최근 30일 가격 추세 강도",
            ],
        },
        {
            "page": 7,
            "title": "시그널 카드 해석 가이드",
            "img": SCREENS / "08-signals.png",
            "paragraphs": [
                "스캐너 결과만으로 매매를 결정하지 마세요. 배지·모멘텀·Magic D는 “관심 목록” 필터이며, "
                "반드시 월가리포트·공시·뉴스와 교차 확인하십시오.",
                "숏 스퀴즈 타겟: 공매도 잔고 대비 거래량·가격 급등 가능성을 모니터링하는 스크리너입니다. "
                "변동성이 크므로 손절 기준을 사전에 정하세요.",
                "고배당 컴파운더: 배당 수익률·현금흐름 안정성 위주. 장기 보유 관점과 단기 모멘텀을 "
                "구분해 읽어야 합니다. Signal Hub에서 다른 이용자의 시그널·카피 트레이드도 참고할 수 있습니다(Day+).",
            ],
        },
        {
            "page": 8,
            "title": "AI Pilot — 대화형 분석",
            "img": PDF_PAGES / "ir-page-03.png",
            "paragraphs": [
                "AI Pilot(Premium+)은 AlphaQuant 대화 인터페이스입니다. "
                "리버스 엔지니어링 종목 선정, 촉매 순위, 진입·손절·목표가 가이드를 자연어로 요청할 수 있습니다.",
                "상단 칩 버튼은 자주 쓰는 프롬프트 예시입니다(예: 역발상 현금흐름 국내 5종, 단기 성장 잠재력 순). "
                "「최근 스캐너 결과와 연동됨」 표시 시, 직전 스캔 컨텍스트가 AI 답변에 반영됩니다.",
                "종목명·티커를 입력하면(예: Robinhood HOOD 주가분석) AI가 가격·PER·ROE·사업모델·촉매를 "
                "구조화해 답합니다. AI 답변은 사실 확인용이며, 공시 원문과 반드시 대조하세요.",
            ],
            "bullets": [
                "4 Pillars: Neglected Bottom + Cash Fortress + Megatrend + 4주 내 촉매",
                "WallPilot AI(서버) 또는 로컬 모델 — 상태 배지로 확인",
                "투자 권유가 아닌 참고 해석",
            ],
        },
        {
            "page": 9,
            "title": "AI Pilot — 추천 종목·매매 가이드",
            "img": PDF_PAGES / "ir-page-04.png",
            "paragraphs": [
                "「스캐너 추천 종목 · 추천 이유」 패널은 스캔 결과 상위 종목에 대해 실행 가능한 "
                "가격 구간을 제시합니다. 각 카드는 세 영역으로 구성됩니다.",
                "① 헤더: 회사명·티커·현재가·스크리너 태그(숏 스퀴즈 등)·OVERWEIGHT/HOLD.",
                "② 가격 전략 박스: 매수 구간(Buy Range), 익절 목표(Take Profit), 하드 손절(Hard Stop). "
                "분할 매수·지정가 주문 계획 수립에 참고하세요.",
                "③ 추천 이유: 스크리너 순위, 30일 모멘텀, Magic D, RSI·MACD·볼린저 %B, Bull/Bear 요약, "
                "거래량 대비 20일 평균, 리스크 메모. 「Balanced view — maintain Overweight」 등 "
                "판정 문구는 AI·퀀트 모델의 종합 의견입니다.",
            ],
        },
        {
            "page": 10,
            "title": "월가리포트 — 검색·퀀트 모델",
            "img": PDF_PAGES / "ir-page-05.png",
            "paragraphs": [
                "Wall St. Report(Day+)는 Peter Lynch GARP, Joel Greenblatt Magic Formula, "
                "한국 수급(외국인·기관)을 한 화면에 통합한 종목 분석 리포트입니다.",
                "검색창에 6자리 국내 코드(005930), 미국 티커(NVDA), 회사명을 입력 후 "
                "「리포트 생성」을 클릭합니다. 「에이전트 심층 리포트」는 AI 에이전트 기반 장문 분석(Premium).",
                "요약 카드: 적정가(Fair Value), 안전마진(Margin of Safety), 매수 구간, 익절 목표. "
                "Lynch 패널: PEG·Lynch Score·Upside. Greenblatt: ROIC·Earnings Yield·Magic Grade. "
                "수급 패널: MCP/증권 API 연동 시 외국인·기관 순매수 표시.",
            ],
            "bullets": [
                "HOLD/OVERWEIGHT/UNDERWEIGHT — 모델 종합 등급",
                "RSI·MACD·BB %B — 기술적 지표 한 줄 요약",
                "실시간 차트: 약 45초 간격 갱신",
            ],
        },
        {
            "page": 11,
            "title": "월가리포트 — 차트·심층 해설",
            "img": PDF_PAGES / "ir-page-06.png",
            "paragraphs": [
                "가격 차트는 3개월 전후 추세를 영역 채우기 라인으로 표시합니다. "
                "현재가·등락률·USD 환산(1 USD = ₩1,512 등)을 함께 확인하세요.",
                "「예정된 촉매」: 30D momentum, 거래량 vs 20일 평균, 거시 리스크 메모.",
                "「에이전트 심층 리포트」: Bull vs Bear, Verdict(판정), Risk Gate(리스크 통과 여부). "
                "한국어 상세 해설에는 Research Manager 관점, 공격적/보수적 실행 가이드, "
                "매수 구간·익절·손절 기준이 구체적으로 기재됩니다.",
                "하단 면책: 연구 목적 정보이며 투자자문이 아님. KR 수급은 MCP, US는 Yahoo Finance 등 출처 명시.",
            ],
        },
        {
            "page": 12,
            "title": "DARTLAB·시그널·토스 트레이더",
            "img": SCREENS / "05-dartlab.png",
            "paragraphs": [
                "DARTLAB(Day+): OpenDART 공시·재무 데이터를 WallPilot UI에서 조회·분석. "
                "PDF 내보내기 지원. AI 해석 전 원문 공시 확인용으로 활용하세요.",
                "Signal Hub(Day+): 커뮤니티 시그널 게시·카피. Live Signals와 연계해 아이디어를 공유합니다.",
                "Toss Trader(Day+, 실행은 Elite): 토스증권 Open API 연동 패널. "
                "잔고·보유 종목·주문 UI. API 키는 My API에서 등록 후 상단 「토스 API 연결」 상태가 녹색으로 바뀝니다.",
                "Agent Desk(Premium): 멀티 에이전트 리서치 데스크. RL Lab(Elite): 강화학습 실험 환경.",
            ],
            "bullets": [
                "dartlab — /dartlab",
                "signals — /signals",
                "toss-trader — /toss-trader",
                "agents/desk — /agents/desk",
            ],
        },
        {
            "page": 13,
            "title": "My API·요금제 관리",
            "img": SCREENS / "09-my-api.png",
            "paragraphs": [
                "My API(Free 열람, Day+ 실행): 토스 Open API App Key·Secret 등록, "
                "연결 테스트, Webhook URL(해당 시) 설정. 키는 서버에 암호화 저장되며 RLS로 본인만 접근.",
                "Pricing: 현재 등급·다음 결제일·Stripe Customer Portal 링크. "
                "플랜 업/다운그레이드는 결제 주기 종료 시점 또는 즉시(정책에 따름) 반영.",
                "Elite 전용 토스 주문 실행 전 반드시 모의·소액으로 API 권한(scope)과 "
                "주문 유형(지정가/시장가)을 확인하세요.",
            ],
        },
        {
            "page": 14,
            "title": "FAQ·안전 수칙",
            "img": None,
            "paragraphs": [
                "Q. WallPilot은 투자자문인가요? — 아닙니다. 참고용 분석 도구이며 손실 책임은 이용자에게 있습니다.",
                "Q. 매수·매도 신호를 주나요? — 적정가·관심 구간·리스크 시나리오를 제공할 뿐, "
                "일률적 매매 지시가 아닙니다.",
                "Q. AI 답변을 믿어도 되나요? — 보조 도구입니다. DART·증권사 API 원문과 교차 확인하세요.",
                "Q. 한국·미국 모두 되나요? — 스캔·리포트는 KR/US 지원. DART·토스는 한국 상장·토스 연동 시.",
            ],
            "bullets": [
                "분산 투자·손절 원칙을 스스로 정하고 준수",
                "레버리지·전 재산 투입 금지",
                "공시·실적 발표 전후 변동성 주의",
                "API 키·계정 공유 금지",
            ],
        },
        {
            "page": 15,
            "title": "법적 고지·지원",
            "img": None,
            "paragraphs": [
                "WallPilot Pro™는 Terrabridge Capital Inc.의 독점 소프트웨어입니다. "
                "Inventor / IP 문의: kangjunchul8@gmail.com",
                "본 서비스는 금융투자업법상 투자권유·투자자문·투자일임을 목적으로 하지 않습니다. "
                "과거 데이터·모델 결과는 미래 수익을 보장하지 않습니다.",
                "이용약관: /terms · 개인정보처리방침: /privacy · "
                "기술·결제 문의: shinkang888@gmail.com (운영)",
                "문서 버전: v1.0 · 스크린샷 출처: wallpilotir.pdf 및 wallpilotpro.vercel.app",
            ],
        },
    ]

    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)

    add_title_slide(
        prs,
        "WallPilot Pro\n이용자 매뉴얼",
        "Reverse-Quant · KR · US · 15페이지 상세 가이드",
        f"Terrabridge Capital Inc. · {date.today().year} · v1.0",
    )

    for p in pages:
        add_manual_slide(
            prs,
            page=p["page"],
            title=p["title"],
            paragraphs=p["paragraphs"],
            img=p.get("img"),
            bullets=p.get("bullets"),
        )

    OUT_MD.write_text(build_markdown(pages), encoding="utf-8")
    prs.save(str(OUT_PPTX))
    print(f"PPTX: {OUT_PPTX}")
    print(f"MD:   {OUT_MD}")
    print(f"Slides: {len(prs.slides)}")


if __name__ == "__main__":
    main()
