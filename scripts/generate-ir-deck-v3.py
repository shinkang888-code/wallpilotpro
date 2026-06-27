#!/usr/bin/env python3
"""WallPilot Pro IR Deck v3 — 40 slides, image-first, data-driven."""
from __future__ import annotations

import importlib.util
import shutil
from datetime import date
from pathlib import Path

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt

ROOT = Path(__file__).resolve().parents[1]
_v2_path = ROOT / "scripts" / "generate-ir-deck.py"
_spec = importlib.util.spec_from_file_location("generate_ir_deck", _v2_path)
_ir_v2 = importlib.util.module_from_spec(_spec)
assert _spec.loader is not None
_spec.loader.exec_module(_ir_v2)

chart_competitor_pricing = _ir_v2.chart_competitor_pricing
chart_market_tam = _ir_v2.chart_market_tam
chart_mrr_scenarios = _ir_v2.chart_mrr_scenarios
chart_pnl = _ir_v2.chart_pnl
chart_revenue_streams = _ir_v2.chart_revenue_streams
chart_site_structure = _ir_v2.chart_site_structure
chart_tier_mix = _ir_v2.chart_tier_mix
chart_unit_economics = _ir_v2.chart_unit_economics

ASSETS = ROOT / "assets"
PDF_PAGES = ASSETS / "ir-pdf-pages"
SCREENS = ASSETS / "ir-screens"
SECTIONS = ASSETS
OUT_DIR = Path.home() / "Downloads" / f"WallPilot_Pro_IR_{date.today().isoformat()}_v3"
OUT_PPTX = OUT_DIR / "WallPilot_Pro_IR_Business_Plan_40slides_v3.pptx"

# Brand — matches live UI (dark nav + gold accent + light content)
NAVY = RGBColor(0x0A, 0x0A, 0x0A)
NAVY_SOFT = RGBColor(0x11, 0x18, 0x27)
GOLD = RGBColor(0xD4, 0xAF, 0x37)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
OFF_WHITE = RGBColor(0xF8, 0xFA, 0xFC)
SLATE = RGBColor(0x64, 0x74, 0x8B)
GRAY = RGBColor(0x9C, 0xA3, 0xAF)
CYAN = RGBColor(0x22, 0xD3, 0xEE)
GREEN = RGBColor(0x34, 0xD3, 0x99)
RED = RGBColor(0xF8, 0x71, 0x71)

W = Inches(13.333)
H = Inches(7.5)


def ensure() -> None:
    ASSETS.mkdir(parents=True, exist_ok=True)
    OUT_DIR.mkdir(parents=True, exist_ok=True)


def pdf(n: int) -> Path:
    return PDF_PAGES / f"ir-page-{n:02d}.png"


def screen(name: str) -> Path:
    return SCREENS / f"{name}.png"


def set_bg(slide, color: RGBColor = NAVY) -> None:
    slide.background.fill.solid()
    slide.background.fill.fore_color.rgb = color


def footer(slide, text: str = "Terrabridge Capital Inc. · Confidential · wallpilotpro.vercel.app") -> None:
    box = slide.shapes.add_textbox(Inches(0.45), Inches(7.08), Inches(12.4), Inches(0.32))
    p = box.text_frame.paragraphs[0]
    p.text = text
    p.font.size = Pt(8)
    p.font.color.rgb = GRAY


def gold_line(slide, top: float = 0.92) -> None:
    ln = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.45), Inches(top), Inches(12.4), Inches(0.04))
    ln.fill.solid()
    ln.fill.fore_color.rgb = GOLD
    ln.line.fill.background()


def title_text(slide, text: str, top: float = 0.38, size: int = 34, color: RGBColor = GOLD, align=PP_ALIGN.LEFT) -> None:
    tb = slide.shapes.add_textbox(Inches(0.55), Inches(top), Inches(12.2), Inches(0.9))
    p = tb.text_frame.paragraphs[0]
    p.text = text
    p.font.size = Pt(size)
    p.font.bold = True
    p.font.color.rgb = color
    p.alignment = align


def caption_bar(slide, title: str, subtitle: str = "", light: bool = False) -> None:
    """Bottom caption strip over full-bleed images."""
    bar_h = Inches(1.35)
    top = H - bar_h - Inches(0.08)
    if light:
        bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0), top, W, bar_h + Inches(0.08))
        bar.fill.solid()
        bar.fill.fore_color.rgb = WHITE
        try:
            bar.fill.transparency = 0.06
        except Exception:
            pass
        bar.line.fill.background()
        tcolor, scolor = NAVY, SLATE
    else:
        bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0), top, W, bar_h + Inches(0.08))
        bar.fill.solid()
        bar.fill.fore_color.rgb = NAVY
        try:
            bar.fill.transparency = 0.18
        except Exception:
            pass
        bar.line.fill.background()
        tcolor, scolor = GOLD, GRAY
    tb = slide.shapes.add_textbox(Inches(0.65), top + Inches(0.18), Inches(12), Inches(1.0))
    tf = tb.text_frame
    p = tf.paragraphs[0]
    p.text = title
    p.font.size = Pt(22)
    p.font.bold = True
    p.font.color.rgb = tcolor
    if subtitle:
        p2 = tf.add_paragraph()
        p2.text = subtitle
        p2.font.size = Pt(12)
        p2.font.color.rgb = scolor
        p2.space_before = Pt(4)


def hero_image(slide, img: Path, overlay: float = 0.0) -> None:
    if img.exists():
        slide.shapes.add_picture(str(img), Inches(0), Inches(0), width=W, height=H)
    if overlay > 0:
        ov = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), W, H)
        ov.fill.solid()
        ov.fill.fore_color.rgb = NAVY
        ov.line.fill.background()
        try:
            ov.fill.transparency = overlay
        except Exception:
            pass


def product_hero(slide, img: Path, title: str, subtitle: str) -> None:
    """Light product slide — image dominant, app-like white frame."""
    set_bg(slide, OFF_WHITE)
    frame = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.35), Inches(0.55), Inches(12.65), Inches(6.35))
    frame.fill.solid()
    frame.fill.fore_color.rgb = WHITE
    frame.line.color.rgb = RGBColor(0xE2, 0xE8, 0xF0)
    frame.line.width = Pt(1)
    if img.exists():
        slide.shapes.add_picture(str(img), Inches(0.45), Inches(0.65), width=Inches(12.45), height=Inches(5.55))
    caption_bar(slide, title, subtitle, light=True)
    footer(slide, "WallPilot Pro™ · Live Product Screenshot · wallpilotir.pdf")


def full_bleed_product(slide, img: Path, title: str, subtitle: str) -> None:
    set_bg(slide, NAVY)
    hero_image(slide, img)
    caption_bar(slide, title, subtitle, light=False)
    footer(slide)


def section_slide(prs: Presentation, bg: Path, section: str, tagline: str) -> None:
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(slide)
    hero_image(slide, bg, overlay=0.35)
    title_text(slide, section, top=2.8, size=44, align=PP_ALIGN.CENTER)
    tb = slide.shapes.add_textbox(Inches(1.5), Inches(3.75), Inches(10.3), Inches(0.8))
    p = tb.text_frame.paragraphs[0]
    p.text = tagline
    p.font.size = Pt(16)
    p.font.color.rgb = CYAN
    p.alignment = PP_ALIGN.CENTER
    gold_line(slide, top=3.55)
    footer(slide)


def stat_cards(slide, stats: list[tuple[str, str, str]]) -> None:
    """label, value, note"""
    n = len(stats)
    gap = 12.2 / n
    for i, (label, value, note) in enumerate(stats):
        left = 0.55 + i * gap
        card = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(left), Inches(1.35), Inches(gap - 0.15), Inches(1.55))
        card.fill.solid()
        card.fill.fore_color.rgb = NAVY_SOFT
        card.line.color.rgb = GOLD
        card.line.width = Pt(0.75)
        tb = slide.shapes.add_textbox(Inches(left + 0.12), Inches(1.48), Inches(gap - 0.35), Inches(1.3))
        tf = tb.text_frame
        p = tf.paragraphs[0]
        p.text = label
        p.font.size = Pt(10)
        p.font.color.rgb = GRAY
        p2 = tf.add_paragraph()
        p2.text = value
        p2.font.size = Pt(26)
        p2.font.bold = True
        p2.font.color.rgb = GOLD
        p2.space_before = Pt(2)
        p3 = tf.add_paragraph()
        p3.text = note
        p3.font.size = Pt(9)
        p3.font.color.rgb = CYAN
        p3.space_before = Pt(4)


def bullets(slide, items: list[str], left: float = 0.55, top: float = 1.2, width: float = 12.2, size: int = 14) -> None:
    box = slide.shapes.add_textbox(Inches(left), Inches(top), Inches(width), Inches(5.8))
    tf = box.text_frame
    tf.word_wrap = True
    for i, item in enumerate(items):
        para = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        para.text = item if item.startswith("•") else f"• {item}"
        para.font.size = Pt(size)
        para.font.color.rgb = WHITE
        para.space_after = Pt(8)


def table_slide(slide, title: str, headers: list[str], rows: list[list[str]]) -> None:
    set_bg(slide)
    title_text(slide, title, top=0.32, size=26)
    gold_line(slide, top=0.88)
    cols = len(headers)
    rws = len(rows) + 1
    th = min(0.38 * rws, 5.2)
    tbl = slide.shapes.add_table(rws, cols, Inches(0.5), Inches(1.05), Inches(12.35), Inches(th)).table
    for j, h in enumerate(headers):
        c = tbl.cell(0, j)
        c.text = h
        c.fill.solid()
        c.fill.fore_color.rgb = RGBColor(0x1F, 0x29, 0x37)
        for p in c.text_frame.paragraphs:
            p.font.bold = True
            p.font.size = Pt(11)
            p.font.color.rgb = GOLD
    for i, row in enumerate(rows, start=1):
        for j, val in enumerate(row):
            c = tbl.cell(i, j)
            c.text = val
            c.fill.solid()
            c.fill.fore_color.rgb = NAVY_SOFT if i % 2 else RGBColor(0x0F, 0x17, 0x2A)
            for p in c.text_frame.paragraphs:
                p.font.size = Pt(10)
                p.font.color.rgb = WHITE
    footer(slide)


def chart_slide(slide, title: str, chart: Path, notes: list[str] | None = None) -> None:
    set_bg(slide)
    title_text(slide, title, top=0.32, size=24)
    if chart.exists():
        slide.shapes.add_picture(str(chart), Inches(0.35), Inches(0.95), width=Inches(12.65))
    if notes:
        nb = slide.shapes.add_textbox(Inches(0.55), Inches(6.55), Inches(12.2), Inches(0.75))
        tf = nb.text_frame
        for i, n in enumerate(notes):
            para = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
            para.text = f"▸ {n}"
            para.font.size = Pt(10)
            para.font.color.rgb = GRAY
    footer(slide)


def cover_slide(prs: Presentation) -> None:
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(slide)
    hero_image(slide, pdf(1), overlay=0.42)
    tb = slide.shapes.add_textbox(Inches(0.75), Inches(1.6), Inches(11.5), Inches(2.8))
    tf = tb.text_frame
    p = tf.paragraphs[0]
    p.text = "WallPilot Pro™"
    p.font.size = Pt(52)
    p.font.bold = True
    p.font.color.rgb = GOLD
    p2 = tf.add_paragraph()
    p2.text = "Investor Relations · Business Plan 2026"
    p2.font.size = Pt(22)
    p2.font.color.rgb = WHITE
    p2.space_before = Pt(10)
    p3 = tf.add_paragraph()
    p3.text = "Reverse-Quant · KR · US · Data-Driven One Screen OS"
    p3.font.size = Pt(14)
    p3.font.color.rgb = CYAN
    p3.space_before = Pt(14)
    stat_cards(slide, [
        ("M3 MRR (Base)", "₩4.6M", "80 paid · ARR ₩55M"),
        ("Blended ARPU", "₩56.2K", "Day 60% · Prem 30%"),
        ("LTV / CAC", "8.4×", "Payback ~2 mo"),
        ("Live URL", "v3", "wallpilotpro.vercel.app"),
    ])
    footer(slide, f"Terrabridge Capital Inc. · {date.today().isoformat()} · IR v3 · Confidential")


def build() -> Path:
    ensure()
    charts = {
        "mrr": chart_mrr_scenarios(),
        "tier": chart_tier_mix(),
        "rev": chart_revenue_streams(),
        "pnl": chart_pnl(),
        "tam": chart_market_tam(),
        "unit": chart_unit_economics(),
        "comp": chart_competitor_pricing(),
        "sitemap": chart_site_structure(),
    }
    sec_product = SECTIONS / "ir-v3-section-product.png"
    sec_finance = SECTIONS / "ir-v3-section-finance.png"
    sec_gtm = SECTIONS / "ir-v3-section-gtm.png"

    prs = Presentation()
    prs.slide_width = W
    prs.slide_height = H

    # ── 1 Cover ──
    cover_slide(prs)

    # ── 2 Disclaimer ──
    s = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(s)
    title_text(s, "Important Disclaimer", size=28)
    gold_line(s)
    bullets(s, [
        "본 자료는 Terrabridge Capital Inc.의 WallPilot Pro™ IR·사업 설명용입니다.",
        "투자자문·투자일임·집합투자업이 아니며, 증권 매매 권유가 아닙니다.",
        "경쟁사·시장·MRR 데이터는 GTM 시나리오·공개 요금 기반 reference 추정치입니다.",
        "Proprietary IP — 무단 복제·배포·클론·스크래핑 금지 (IP Shield enforced).",
    ], top=1.15, size=15)
    footer(s)

    # ── 3 Section Vision ──
    section_slide(prs, sec_product, "01 · Vision & Opportunity", "KR+US Reverse-Quant SaaS — One Screen OS")

    # ── 4 Executive Summary ──
    s = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(s)
    title_text(s, "Executive Summary", size=28)
    gold_line(s)
    stat_cards(s, [
        ("Product", "10 Menus", "Scanner → Toss → RL"),
        ("Pricing", "₩0–199K", "4-Tier B2C + B2B"),
        ("M3 MRR", "₩4.6M", "Base · 80 paid"),
        ("Regulatory", "Not IA", "Reference tool only"),
    ])
    if screen("02-scanner").exists():
        s.shapes.add_picture(str(screen("02-scanner")), Inches(0.55), Inches(3.05), width=Inches(12.25))
    footer(s)

    # ── 5 North Star ──
    s = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(s)
    title_text(s, "North Star Metric", size=28)
    gold_line(s)
    stat_cards(s, [
        ("Primary KPI", "WAU Scan", "Paid user ≥2 scans/week"),
        ("Free→Paid CR", "8–12%", "90-day target"),
        ("Day→Premium", "15%", "Upgrade of Day base"),
        ("Churn (Base)", "10%/mo", "Conservative 15%"),
    ])
    bullets(s, [
        "Google OAuth → Free/active 자동 부여 · Stripe/Danal 월구독으로 승급",
        "결제 실패·취소 시 Free(inactive) 자동 다운그레이드",
        "Admin: shinkang888@gmail.com · Product: kangjunchul8@gmail.com",
    ], top=3.2, size=13)
    footer(s)

    # ── 6 Problem ──
    s = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(s)
    title_text(s, "Problem — 정보 격차 & 도구 분산", size=26)
    if pdf(2).exists():
        s.shapes.add_picture(str(pdf(2)), Inches(6.85), Inches(0.95), width=Inches(6.15))
    bullets(s, [
        "차트·재무·수급·뉴스·13F가 앱·탭·커뮤니티에 분산 → 의사결정 지연",
        "퀀터스·스노우볼: 백테스트/ETF 중심 — KR/US 개별주+수급 통합 부족",
        "로보어드바이저(핀트·토스 AI): AUM % — DIY 분석 니즈와 다름",
        "Seeking Alpha·TradingView: US/KR·DART·토스 실행 One OS 아님",
    ], top=1.05, width=6.1, size=13)
    footer(s)

    # ── 7 Problem Data ──
    s = prs.slides.add_slide(prs.slide_layouts[6])
    table_slide(s, "Market Gap — Retail Investor Pain Points", ["Pain", "Impact", "Incumbent Gap"], [
        ["정보 분산", "평균 4+ 앱 전환", "통합 OS 없음"],
        ["KR·US 교차", "환율·수급 이중 추적", "단일 시장 중심"],
        ["13F·하이롤러", "지연·해석 난이도", "원시 데이터만 제공"],
        ["AI 해석", "환각·면책 리스크", "단순 챗봇 수준"],
        ["실행 단절", "분석→주문 friction", "브로커 앱 별도"],
    ])

    # ── 8 Solution ──
    s = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(s)
    title_text(s, "Solution — WallPilot One Screen OS", size=26)
    gold_line(s)
    stat_cards(s, [
        ("① Scan", "Reverse-Quant", "13F·Volume·Quant"),
        ("② Analyze", "Wall Report", "Lynch·Greenblatt"),
        ("③ Brief", "DART·AI", "공시·Gemini Q&A"),
        ("④ Execute", "Toss API", "Elite orders"),
    ])
    bullets(s, [
        "Scanner → Report → DARTLAB → AI Pilot → Agent Desk → Signal Hub → Toss → RL Lab",
        "8개 언어 i18n · PWA · Supabase RLS · Vercel edge deploy",
    ], top=3.15, size=13)
    footer(s)

    # ── 9 Site Structure ──
    s = prs.slides.add_slide(prs.slide_layouts[6])
    chart_slide(s, "Product Architecture — Site Map & Tier Gates", charts["sitemap"], [
        "Free: Scanner preview · Signal Hub read · My API · Pricing",
        "Day+: Report · DART · Toss dashboard · Signals write",
        "Premium+: AI Pilot · Agent Desk · PDF export",
        "Elite: RL Lab · Toss execute (toss_execute entitlement)",
    ])

    # ── 10 Section Product ──
    section_slide(prs, sec_product, "02 · Live Product", "Screenshots from wallpilotir.pdf · Production UI")

    # ── 11 Reverse Quant Engine ──
    s = prs.slides.add_slide(prs.slide_layouts[6])
    product_hero(s, pdf(1), "Reverse-Quant Engine", "상위 20 구루 트래커 · 13F 고래 · NYSE/NASDAQ · KOSPI/KOSDAQ")

    # ── 12 Live Signals ──
    s = prs.slides.add_slide(prs.slide_layouts[6])
    full_bleed_product(s, pdf(2), "Live Signals — Real-time Screener", "Short Squeeze · High Dividend · OVERWEIGHT/HOLD/UNDERWEIGHT · Magic D · 30D Momentum")

    # ── 13 Scanner Metrics ──
    s = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(s)
    title_text(s, "Scanner — Data Points per Signal Card", size=24)
    if screen("02-scanner").exists():
        s.shapes.add_picture(str(screen("02-scanner")), Inches(0.45), Inches(0.95), width=Inches(7.5))
    stat_cards(s, [
        ("Universe", "KR+US", "KOSPI·KOSDAQ·NYSE·NASDAQ"),
        ("Screeners", "2+", "Short Squeeze · Dividend"),
        ("Signals", "20+", "Live matrix columns"),
        ("Refresh", "~45s", "Price · Volume"),
    ])
    bullets(s, [
        "Toss Open API + Yahoo Finance + proprietary quant engine",
        "Free: preview · Day+: full scan execute + chart",
    ], left=8.15, top=3.2, width=4.8, size=11)
    footer(s)

    # ── 14 Landing GTM ──
    s = prs.slides.add_slide(prs.slide_layouts[6])
    product_hero(s, screen("01-home-landing"), "Landing & GTM Hero", "3 Value Props · FAQ 5 · Terms/Privacy · CTA: 무료 스캐너 미리보기 · @vercel/analytics")

    # ── 15 Wall Report Search ──
    s = prs.slides.add_slide(prs.slide_layouts[6])
    full_bleed_product(s, pdf(5), "Wall St. Report — Quant Valuation", "Fair Value ₩208,863 · Lynch PEG · Greenblatt ROIC · RSI/MACD/BB · Buy Range · Take Profit")

    # ── 16 Wall Report Deep ──
    s = prs.slides.add_slide(prs.slide_layouts[6])
    full_bleed_product(s, pdf(6), "Agent Deep Report — Bull/Bear · Risk Gate", "3M chart · Catalyst · KR execution guide · MCP/Yahoo data sources")

    # ── 17 AI Pilot Chat ──
    s = prs.slides.add_slide(prs.slide_layouts[6])
    full_bleed_product(s, pdf(3), "AI Pilot — AlphaQuant Conversation", "4 Pillars · Scanner-linked context · HOOD analysis example · Premium+ tier")

    # ── 18 AI Recommendations ──
    s = prs.slides.add_slide(prs.slide_layouts[6])
    full_bleed_product(s, pdf(4), "AI Pilot — Trade Setup Cards", "Buy Range · Take Profit · Hard Stop · RSI · MACD · Bull/Bear verdict · 20 picks")

    # ── 19 DARTLAB ──
    s = prs.slides.add_slide(prs.slide_layouts[6])
    img = screen("05-dartlab") if screen("05-dartlab").exists() else pdf(5)
    product_hero(s, img, "DARTLAB — OpenDART KR Filings", "재무·공시 AI brief · PDF export · Day+ tier · Cross-check 원문 권장")

    # ── 20 Toss Trader ──
    s = prs.slides.add_slide(prs.slide_layouts[6])
    product_hero(s, screen("07-toss-trader"), "Toss Trader — Open API Dashboard", "보유·손익·미체결·매수가능 · Day+ view · Elite: toss_execute orders")

    # ── 21 Signals + Pricing ──
    s = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(s, OFF_WHITE)
    title_text(s, "Signal Hub & Pricing", size=24, color=NAVY)
    if screen("08-signals").exists():
        s.shapes.add_picture(str(screen("08-signals")), Inches(0.45), Inches(0.95), width=Inches(6.2))
    if screen("03-pricing").exists():
        s.shapes.add_picture(str(screen("03-pricing")), Inches(6.85), Inches(0.95), width=Inches(6.2))
    footer(s, "Signal Hub: Free read · Day+ post/copy · Pricing: Stripe + Danal")

    # ── 22 Tier Matrix ──
    s = prs.slides.add_slide(prs.slide_layouts[6])
    table_slide(s, "Entitlement Matrix — Menu × Tier", ["Menu", "Free", "Day", "Premium", "Elite"], [
        ["Scanner", "preview", "execute", "execute+PDF", "execute+PDF"],
        ["Wall Report", "—", "execute", "execute+PDF", "execute+PDF"],
        ["DARTLAB", "view", "execute+PDF", "execute+PDF", "execute+PDF"],
        ["AI Pilot", "—", "—", "execute+PDF", "execute+PDF"],
        ["Signal Hub", "view", "execute", "execute+PDF", "execute+PDF"],
        ["Toss Trader", "view", "dashboard", "dashboard", "execute orders"],
        ["RL Lab", "—", "—", "—", "execute+PDF"],
    ])

    # ── 23 Section Business ──
    section_slide(prs, sec_finance, "03 · Business Model", "4-Tier B2C · B2B API · Unit Economics")

    # ── 24 Pricing Full ──
    s = prs.slides.add_slide(prs.slide_layouts[6])
    product_hero(s, screen("03-pricing"), "Monetization — 4-Tier Subscription", "Free ₩0 · Day ₩39K · Premium ₩99K · Elite ₩199K · US: $0/$29/$59/$99")

    # ── 25 Subscription Flow ──
    s = prs.slides.add_slide(prs.slide_layouts[6])
    table_slide(s, "Subscription System — Auth & Billing Flow", ["Event", "DB plan", "Status"], [
        ["Google OAuth signup", "free", "active"],
        ["Stripe/Danal paid", "pro/premium/elite", "active/trialing"],
        ["Payment failed/canceled", "free", "inactive"],
        ["Admin override", "manual tier", "active"],
        ["Webhook", "STRIPE_PRICE_* env", "auto sync"],
    ])

    # ── 26 Competitive Table ──
    s = prs.slides.add_slide(prs.slide_layouts[6])
    table_slide(s, "Competitive Landscape — 2026 Benchmark", ["Player", "Model", "Price", "vs WallPilot"], [
        ["WallPilot Pro", "Flat SaaS", "₩39~199K/mo", "KR+US quant+AI+Toss OS"],
        ["퀀터스", "구독+실전권", "연 ₩199K~359K", "백테스트·자동매매"],
        ["스노우볼72", "Freemium", "₩0", "ETF only"],
        ["Seeking Alpha", "Flat SaaS", "~$299/yr", "US · no DART"],
        ["TradingView", "Flat SaaS", "$12~60/mo", "Charting"],
        ["핀트/토스 AI", "AUM %", "0.55~1.17%/yr", "일임·로보"],
        ["Bloomberg", "Terminal", "~$2K/mo", "Institutional"],
    ])

    # ── 27 Competitor Pricing Chart ──
    s = prs.slides.add_slide(prs.slide_layouts[6])
    chart_slide(s, "Pricing Benchmark — Monthly KRW Equivalent", charts["comp"], [
        "WallPilot Day ₩39K vs SA Premium ₩34K · TV Pro ₩18K",
        "Premium ₩99K includes AI Agent + PDF — differentiated bundle",
    ])

    # ── 28 B2B ──
    s = prs.slides.add_slide(prs.slide_layouts[6])
    table_slide(s, "B2B License & API Revenue", ["Tier", "USD/mo", "Includes"], [
        ["Starter", "$2,500", "1K calls/day · 1 brand"],
        ["Growth", "$8,000", "10K calls/day · co-brand · SLA 99.5%"],
        ["Enterprise", "$25,000+", "Unlimited* · white-label"],
        ["PoC (8wk)", "$5,000 flat", "Sandbox · embed · KPI review"],
        ["White-label setup", "$15K~$50K", "One-time integration"],
    ])

    # ── 29 Section Market ──
    section_slide(prs, sec_finance, "04 · Market & Financials", "TAM · MRR Scenarios · 90-Day P&L")

    # ── 30 TAM ──
    s = prs.slides.add_slide(prs.slide_layouts[6])
    chart_slide(s, "Market Opportunity — TAM / SAM / SOM", charts["tam"], [
        "TAM ~$120B global retail quant·AI",
        "SAM ~$18B KR+US active traders",
        "SOM ~$2.4B WallPilot Y3 target",
    ])

    # ── 31 MRR ──
    s = prs.slides.add_slide(prs.slide_layouts[6])
    chart_slide(s, "MRR Projection — 3 Scenarios (₩M)", charts["mrr"], [
        "M3 Base: ₩4.6M (64 KR + 16 US paid)",
        "M12 Base: ~₩28M MRR · ARR run-rate ~₩336M",
    ])

    # ── 32 Tier Mix ──
    s = prs.slides.add_slide(prs.slide_layouts[6])
    chart_slide(s, "Revenue Model — Tier Mix & Blended ARPU", charts["tier"], [
        "Day 60% · Premium 30% · Elite 10%",
        "Blended ARPU ₩56.2K/mo (KR) · ~$45/mo (US)",
    ])

    # ── 33 Revenue Streams ──
    s = prs.slides.add_slide(prs.slide_layouts[6])
    chart_slide(s, "Revenue Streams — B2C · B2B · US", charts["rev"], [
        "M3: B2C ₩4.0M + PoC ₩0.675M + US Stripe ₩0.97M",
        "M12: B2C ₩22M + B2B API ₩8M + US ₩6M",
    ])

    # ── 34 Unit Economics ──
    s = prs.slides.add_slide(prs.slide_layouts[6])
    chart_slide(s, "Unit Economics — CAC · LTV · Payback", charts["unit"], [
        "Blended CAC ≤ ₩40K · LTV(6mo) ₩337K · LTV/CAC 8.4×",
        "Payback ~2 months on Day tier",
    ])

    # ── 35 P&L ──
    s = prs.slides.add_slide(prs.slide_layouts[6])
    chart_slide(s, "90-Day P&L Sketch (Base Scenario)", charts["pnl"], [
        "M3 net positive ~₩2.1M (excl. founder salary)",
        "Marketing M1-M3: ₩1.0→₩2.0M · Infra: ₩0.3→₩0.5M",
    ])

    # ── 36 GTM KPI ──
    s = prs.slides.add_slide(prs.slide_layouts[6])
    table_slide(s, "90-Day GTM KPI — Base Scenario", ["Metric", "Target", "Phase"], [
        ["Free signups", "800", "W1-12 cumulative"],
        ["Paid (M3)", "80 · MRR ₩4.6M", "W9-12"],
        ["Free→Paid CR", "8–12%", "Ongoing"],
        ["WAU / Paid", ">70%", "North star"],
        ["B2B PoC", "1 deal", "W9-12"],
        ["SNS posts", "3/week", "W3-8"],
        ["CAC blended", "≤₩40K", "All channels"],
    ])

    # ── 37 GTM Phases ──
    s = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(s)
    title_text(s, "GTM Roadmap — 90-Day Phases", size=26)
    gold_line(s)
    stat_cards(s, [
        ("Phase 1", "W1-2", "Landing·Terms·Analytics"),
        ("Phase 2", "W3-4", "SNS · Free→Day funnel"),
        ("Phase 3", "W5-8", "Case studies · AI demo"),
        ("Phase 4", "W9-12", "B2B 10 · US EN · Elite beta"),
    ])
    bullets(s, [
        "Phase 1 ✓ Hero · FAQ · /terms · /privacy · Vercel Analytics",
        "Cold email targets: 10 (MTS · Toss API · YouTube research · RIA · IB)",
        "Day 91-180: 250+ paid · MRR ₩12M+ · US revenue 30%+",
    ], top=3.15, size=12)
    footer(s)

    # ── 38 Tech Stack ──
    s = prs.slides.add_slide(prs.slide_layouts[6])
    table_slide(s, "Technical Architecture — Repo Analysis", ["Layer", "Stack", "Role"], [
        ["Frontend", "React 19 · TanStack · PWA · 8 langs", "Landing · Nav · Tier gates"],
        ["Engine", "Quant · Supply · Valuation · Agents", "Reverse-Quant pipeline"],
        ["Data", "OpenDART · 13F · AV · Toss · Gemini", "KR/US cross-market"],
        ["SaaS", "Supabase RLS · Stripe/Danal · Vercel", "Auth · Billing · Deploy"],
        ["Security", "IP Shield · CSP · violation log", "Clone protection"],
    ])

    # ── 39 Roadmap + IP ──
    s = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(s)
    title_text(s, "Roadmap · IP Moat · Risk", size=26)
    gold_line(s)
    if screen("09-my-api").exists():
        s.shapes.add_picture(str(screen("09-my-api")), Inches(7.0), Inches(1.0), width=Inches(5.9))
    bullets(s, [
        "Q2 2026: GTM · 80+ paid · B2B PoC · wallpilotpro.vercel.app live",
        "Q3 2026: US 30% revenue · Enterprise API · GA4 funnel",
        "Q4 2026: 250+ paid · MRR ₩12M+ · Danal full",
        "2027: RIA channel · Mobile native · US trademark",
        "IP: Terrabridge Capital Inc. · WallPilot Pro™ · IP Shield · Terms/Privacy live",
        "Risk: 규제 변경 · API rate · FX · churn · clone (mitigated by IP Shield)",
    ], top=1.05, width=6.2, size=11)
    footer(s)

    # ── 40 Thank You / Ask ──
    s = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(s)
    hero_image(s, pdf(1), overlay=0.5)
    tb = s.shapes.add_textbox(Inches(0.75), Inches(2.0), Inches(11.5), Inches(3.5))
    tf = tb.text_frame
    p = tf.paragraphs[0]
    p.text = "Thank You"
    p.font.size = Pt(48)
    p.font.bold = True
    p.font.color.rgb = GOLD
    p2 = tf.add_paragraph()
    p2.text = "Partnership · PoC · Strategic Investment"
    p2.font.size = Pt(20)
    p2.font.color.rgb = WHITE
    p2.space_before = Pt(12)
    p3 = tf.add_paragraph()
    p3.text = "terrabridgecapital@gmail.com · kangjunchul8@gmail.com"
    p3.font.size = Pt(14)
    p3.font.color.rgb = CYAN
    p3.space_before = Pt(16)
    p4 = tf.add_paragraph()
    p4.text = "wallpilotpro.vercel.app · WallPilot Pro™ IR v3 · Data-Driven Quant Analysis"
    p4.font.size = Pt(12)
    p4.font.color.rgb = GRAY
    p4.space_before = Pt(10)
    footer(s)

    prs.save(str(OUT_PPTX))

    for f in ASSETS.glob("chart_*.png"):
        shutil.copy2(f, OUT_DIR / f.name)
    if PDF_PAGES.is_dir():
        dest = OUT_DIR / "ir-pdf-pages"
        dest.mkdir(exist_ok=True)
        for f in PDF_PAGES.glob("*.png"):
            shutil.copy2(f, dest / f.name)
    if SCREENS.is_dir():
        dest = OUT_DIR / "ir-screens"
        dest.mkdir(exist_ok=True)
        for f in SCREENS.glob("*.png"):
            shutil.copy2(f, dest / f.name)
    for f in SECTIONS.glob("ir-v3-section-*.png"):
        shutil.copy2(f, OUT_DIR / f.name)

    return OUT_PPTX


if __name__ == "__main__":
    path = build()
    print(f"Generated: {path}")
    print(f"Assets: {OUT_DIR}")
    from pptx import Presentation as P

    print(f"Slides: {len(P(str(path)).slides)}")
